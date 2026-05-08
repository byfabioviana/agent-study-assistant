"""
Manual de Uso dos Agentes - Verticalis AI Journey deck
Como o mestrando aciona o squad no dia a dia.
Padrao Oxford Heritage v2.0 | Plus Jakarta Sans + Lora + IBM Plex Mono
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === Oxford Heritage v2.0 tokens ===
PRIMARY_950 = RGBColor(0x0B, 0x28, 0x23)
PRIMARY_900 = RGBColor(0x2D, 0x45, 0x3E)
PRIMARY_700 = RGBColor(0x3A, 0x56, 0x4E)
PRIMARY_600 = RGBColor(0x2D, 0x6E, 0x43)
ACCENT_600 = RGBColor(0xC5, 0xB3, 0x58)
ACCENT_500 = RGBColor(0xC5, 0xB3, 0x58)
ACCENT_400 = RGBColor(0xD4, 0xC8, 0x78)
NEUTRAL_950 = RGBColor(0x1E, 0x1E, 0x1E)
NEUTRAL_800 = RGBColor(0x33, 0x33, 0x33)
NEUTRAL_700 = RGBColor(0x4A, 0x4A, 0x4A)
NEUTRAL_600 = RGBColor(0x66, 0x66, 0x66)
NEUTRAL_500 = RGBColor(0x88, 0x88, 0x88)
NEUTRAL_400 = RGBColor(0xAA, 0xAA, 0xAA)
NEUTRAL_300 = RGBColor(0xCC, 0xCC, 0xCC)
NEUTRAL_200 = RGBColor(0xE0, 0xDD, 0xD5)
NEUTRAL_50 = RGBColor(0xF4, 0xF1, 0xEA)
DEW_GREEN = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xC0, 0x39, 0x2B)
SUCCESS = RGBColor(0x2D, 0x6E, 0x43)
WARN = RGBColor(0xD4, 0x8B, 0x1F)

FONT_HEADING = "Plus Jakarta Sans"
FONT_BODY = "Plus Jakarta Sans"
FONT_DISPLAY = "Lora"
FONT_MONO = "IBM Plex Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ROOT = Path(__file__).resolve().parent.parent

# ----------------------- helpers -----------------------
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, l, t, w, h, text, font=FONT_BODY, sz=18, clr=WHITE,
        bold=False, align=PP_ALIGN.LEFT, spacing=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    if spacing:
        p.line_spacing = Pt(spacing)
    return box, tf


def bar(slide, l, t, w=Inches(1), h=Pt(3), clr=ACCENT_500):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = clr
    s.line.fill.background()
    return s


def left_bar(slide, clr=ACCENT_500):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(6), SLIDE_H)
    s.fill.solid()
    s.fill.fore_color.rgb = clr
    s.line.fill.background()


def section_label(slide, text, l=Inches(1.2), t=Inches(0.6), clr=ACCENT_600):
    txt(slide, l, t, Inches(8), Inches(0.35), text.upper(),
        FONT_HEADING, 11, clr, bold=True)


def chapter_num(slide, num, clr=ACCENT_400):
    txt(slide, Inches(1.2), Inches(0.6), Inches(2), Inches(0.35),
        num, FONT_MONO, 13, clr, bold=True)


def footer(slide, pg, total, dark=True):
    c = NEUTRAL_500 if dark else NEUTRAL_600
    txt(slide, Inches(1.2), Inches(7.05), Inches(8), Inches(0.35),
        "Manual de Uso  |  Agent Study Assistant  |  Verticalis AI Journey",
        FONT_BODY, 9, c)
    txt(slide, Inches(11.0), Inches(7.05), Inches(1.5), Inches(0.35),
        f"{pg:02d} / {total:02d}", FONT_MONO, 9, c, align=PP_ALIGN.RIGHT)


def filled_rect(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = line_w
    return s


def rounded(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = line_w
    return s


def icon_circle(slide, cx, cy, r, fill, mark, mark_color, mark_size=18, mark_font=FONT_HEADING):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - r, cy - r, r*2, r*2)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    txt(slide, cx - r, cy - r + Emu(int(r*0.15)), r*2, r*2,
        mark, mark_font, mark_size, mark_color, bold=True, align=PP_ALIGN.CENTER)


def code_block(slide, l, t, w, h, code, sz=11, clr=ACCENT_400, bg=PRIMARY_950, border=NEUTRAL_700):
    rounded(slide, l, t, w, h, bg, line=border, line_w=Pt(0.5))
    txt(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4), h - Inches(0.3),
        code, FONT_MONO, sz, clr, spacing=14)


# ============================================================================
# BUILD
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
TOTAL = 18
BLANK = prs.slide_layouts[6]


# ---------------- 01 CAPA ----------------
def s01_cover():
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    filled_rect(s, Inches(8.5), Inches(4.5), Inches(4.8), Inches(3), PRIMARY_900)

    txt(s, Inches(1.2), Inches(1.6), Inches(11), Inches(0.5),
        "VERTICALIS AI JOURNEY  /  EDUCATION  /  MANUAL DE USO", FONT_HEADING, 12, ACCENT_400, bold=True)
    bar(s, Inches(1.2), Inches(2.2), Inches(1.2), Pt(3))
    txt(s, Inches(1.2), Inches(2.5), Inches(11), Inches(1.5),
        "Como acionar os agentes\nno dia a dia", FONT_DISPLAY, 46, WHITE, bold=True,
        spacing=54)
    txt(s, Inches(1.2), Inches(4.7), Inches(10), Inches(0.5),
        "Framework 5 Fases  •  Prompts copiaveis  •  Ritmos diario, semanal e por sprint",
        FONT_BODY, 16, NEUTRAL_300)
    txt(s, Inches(1.2), Inches(5.3), Inches(10), Inches(0.4),
        "Da pergunta de pesquisa a defesa, sem voce se perder.",
        FONT_DISPLAY, 17, ACCENT_400)
    txt(s, Inches(1.2), Inches(6.5), Inches(10), Inches(0.4),
        "Maio 2026  •  Confidencial  •  v1.0",
        FONT_MONO, 11, NEUTRAL_500)


# ---------------- 02 INDICE ----------------
def s02_toc(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "INDICE")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Conteudo", FONT_HEADING, 38, NEUTRAL_950, bold=True)
    chapters = [
        ("01", "Pre-requisito", "O que voce precisa ter antes de comecar"),
        ("02", "Anatomia de uma chamada", "Como conversar com um agente"),
        ("03", "Framework 5 FASES", "Mapa do processo do INPUT ao POLIMENTO"),
        ("",   "  Fase 1 — Input",      "academic-researcher + source-validator"),
        ("",   "  Fase 2 — Desenho",    "methodology-advisor"),
        ("",   "  Fase 3 — Digestao",   "knowledge-architect"),
        ("",   "  Fase 4 — Producao",   "academic-writer + citation-manager"),
        ("",   "  Fase 5 — Polimento",  "auditor + humanizer + peer-reviewer"),
        ("04", "Ritmos de uso", "Diario, semanal, sprint, pre-submissao"),
        ("05", "Decision tree", "Qual agente usar quando voce nao sabe"),
        ("06", "Cheat sheet", "Todos os prompts em uma pagina"),
    ]
    for i, (num, title, desc) in enumerate(chapters):
        sub = (num == "")
        y = Inches(2.0 + i * 0.42)
        if num:
            txt(s, Inches(1.2), y, Inches(0.6), Inches(0.4),
                num, FONT_MONO, 13, ACCENT_500, bold=True)
        title_color = NEUTRAL_700 if sub else NEUTRAL_950
        title_size = 13 if sub else 15
        title_bold = False if sub else True
        title_font = FONT_BODY if sub else FONT_HEADING
        txt(s, Inches(2.0), y, Inches(4.5), Inches(0.4),
            title, title_font, title_size, title_color, bold=title_bold)
        desc_color = NEUTRAL_500 if sub else NEUTRAL_600
        desc_size = 11 if sub else 13
        txt(s, Inches(6.5), y, Inches(7), Inches(0.4),
            desc, FONT_BODY, desc_size, desc_color)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 03 MANIFESTO ----------------
def s03_manifesto(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "00")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(1.0),
        '"Os agentes so funcionam se voce souber\nexatamente quando aciona-los."',
        FONT_DISPLAY, 32, WHITE, bold=True, spacing=42)
    bar(s, Inches(1.2), Inches(2.7), Inches(1.2))

    left = (
        "Voce ja instalou os 9 agentes.\n"
        "O ambiente esta pronto.\n"
        "Falta o mais importante: o metodo.\n\n"
        "Sem metodo, voce vai usar o agente errado\n"
        "no momento errado e gastar dois meses\n"
        "tentando reparar uma decisao que custou\n"
        "5 minutos para tomar."
    )
    right = (
        "Este manual oferece um framework simples\n"
        "de 5 fases que voce repete a cada\n"
        "ciclo da pesquisa.\n\n"
        "Cada fase tem agentes especificos,\n"
        "prompts copiaveis, e ritmos de uso\n"
        "(diario, semanal, sprint, pre-submissao)\n"
        "que se encaixam na sua rotina real."
    )
    txt(s, Inches(1.2), Inches(3.2), Inches(5.5), Inches(3.2),
        left, FONT_BODY, 14, NEUTRAL_300, spacing=22)
    txt(s, Inches(7.0), Inches(3.2), Inches(5.5), Inches(3.2),
        right, FONT_BODY, 14, NEUTRAL_300, spacing=22)

    txt(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4),
        "Metodo > Ferramenta. O squad e o instrumento — voce e o pesquisador.",
        FONT_DISPLAY, 16, ACCENT_400, align=PP_ALIGN.CENTER)
    footer(s, pg, TOTAL)


# ---------------- 04 PRE-REQUISITO ----------------
def s04_prerequisito(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "01  PRE-REQUISITO")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Antes de continuar, valide a instalacao",
        FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Se algo aqui nao estiver pronto, volte para o INSTALL-FACIL.md.",
        FONT_BODY, 13, NEUTRAL_600)

    # Checklist
    items = [
        ("OK", "Claude Code instalado",        "App em claude.com/claude-code ou navegador em claude.ai/code"),
        ("OK", "Squad instalado",              "9 agentes em ~/.claude/agents/  (rode /agents para confirmar)"),
        ("OK", "Comando /study-setup ativo",   "Em ~/.claude/commands/study-setup.md"),
        ("OK", "Pasta de trabalho criada",     "Mestrado-{Tema}/ com arvore PARA + Zettelkasten"),
        ("OK", "memory/ pre-preenchido",       "5 arquivos: thesis, citation_style, voice, methodology, feedback"),
        ("?",  "Voice-humanizer calibrado",    "Voce ja respondeu o questionario de 21 perguntas?"),
    ]
    for i, (status, name, desc) in enumerate(items):
        y = Inches(2.6 + i * 0.65)
        rounded(s, Inches(1.2), y, Inches(11.5), Inches(0.55),
                NEUTRAL_50, line=NEUTRAL_200)
        # status badge
        color = SUCCESS if status == "OK" else WARN
        rounded(s, Inches(1.4), y + Inches(0.1), Inches(0.5), Inches(0.35), color)
        txt(s, Inches(1.4), y + Inches(0.13), Inches(0.5), Inches(0.3),
            status, FONT_MONO, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(2.1), y + Inches(0.07), Inches(4), Inches(0.4),
            name, FONT_HEADING, 13, NEUTRAL_950, bold=True)
        txt(s, Inches(6.2), y + Inches(0.13), Inches(6.4), Inches(0.4),
            desc, FONT_BODY, 11, NEUTRAL_600)

    rounded(s, Inches(1.2), Inches(6.55), Inches(11.5), Inches(0.4),
            PRIMARY_950)
    txt(s, Inches(1.4), Inches(6.6), Inches(11.3), Inches(0.3),
        "VALIDAR  •  Cole no Claude Code: /study-setup --check",
        FONT_MONO, 11, ACCENT_400, bold=True)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 05 ANATOMIA DA CHAMADA ----------------
def s05_anatomia(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "02")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Anatomia de uma chamada de agente",
        FONT_HEADING, 30, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "Tres formas de acionar + estrutura ideal de prompt.",
        FONT_BODY, 13, NEUTRAL_400)

    # Tres formas de acionar
    forms = [
        ("01", "Mencao direta", "Acione o {nome-do-agente} para...",
         "Funciona em qualquer\nsessao Claude Code."),
        ("02", "Slash command", "/agents > escolher na lista",
         "Lista os agentes\ninstalados."),
        ("03", "Contexto + tarefa", "Sou {nome}, trabalho com {tema}.\nPara este draft especifico..."),
    ]
    for i, item in enumerate(forms[:3]):
        x = Inches(0.7 + i * 4.2)
        rounded(s, x, Inches(2.85), Inches(4.0), Inches(1.7),
                PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
        txt(s, x + Inches(0.2), Inches(2.95), Inches(0.5), Inches(0.3),
            item[0], FONT_MONO, 11, ACCENT_400, bold=True)
        txt(s, x + Inches(0.85), Inches(2.95), Inches(3), Inches(0.3),
            item[1], FONT_HEADING, 13, WHITE, bold=True)
        code_block(s, x + Inches(0.2), Inches(3.4), Inches(3.7), Inches(0.65),
                    item[2], sz=10)
        if len(item) > 3:
            txt(s, x + Inches(0.2), Inches(4.15), Inches(3.7), Inches(0.4),
                item[3], FONT_BODY, 10, NEUTRAL_300, spacing=14)

    # Estrutura ideal
    rounded(s, Inches(1.2), Inches(4.85), Inches(11.5), Inches(2.0),
            PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
    txt(s, Inches(1.4), Inches(5.0), Inches(11), Inches(0.3),
        "ESTRUTURA IDEAL DE UM PROMPT", FONT_HEADING, 11, ACCENT_500, bold=True)

    parts = [
        ("AGENTE",   "academic-researcher"),
        ("TAREFA",   "buscar fontes recentes"),
        ("CONTEXTO", "tema X, periodo 2020-2026, em SciELO/Scholar"),
        ("SAIDA",    "lista priorizada com DOI verificavel"),
    ]
    for i, (label, content) in enumerate(parts):
        x = Inches(1.4 + (i % 2) * 5.7)
        y = Inches(5.4 + (i // 2) * 0.65)
        txt(s, x, y, Inches(1), Inches(0.3),
            label, FONT_MONO, 10, ACCENT_400, bold=True)
        txt(s, x + Inches(1.2), y, Inches(4), Inches(0.3),
            content, FONT_BODY, 12, WHITE)
    footer(s, pg, TOTAL)


# ---------------- 06 FRAMEWORK OVERVIEW ----------------
def s06_framework(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "03  FRAMEWORK 5 FASES")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Da pergunta a defesa, em ciclos repetiveis",
        FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Cada fase tem agentes especificos, output esperado e ritmo de execucao.",
        FONT_BODY, 13, NEUTRAL_600)

    phases = [
        ("01", "INPUT",     "Capturar e validar fontes",
         "academic-researcher\nsource-validator", "Diario\nSemanal", DANGER),
        ("02", "DESENHO",   "Planejar metodologia",
         "methodology-advisor", "Inicio do projeto\nQualificacao", PRIMARY_700),
        ("03", "DIGESTAO",  "Converter leitura em\nnota atomica",
         "knowledge-architect", "Diario\napos cada leitura", ACCENT_600),
        ("04", "PRODUCAO",  "Redigir capitulos\na partir das notas",
         "academic-writer\ncitation-manager", "Sprints de 1-2 semanas", PRIMARY_950),
        ("05", "POLIMENTO", "Auditar, humanizar,\nsimular banca",
         "ai-pattern-auditor\nvoice-humanizer\npeer-reviewer", "Antes de submeter", SUCCESS),
    ]
    for i, (num, name, desc, agents, when, color) in enumerate(phases):
        x = Inches(0.7 + i * 2.55)
        rounded(s, x, Inches(2.5), Inches(2.4), Inches(4.4),
                NEUTRAL_50, line=NEUTRAL_200)
        bar(s, x, Inches(2.5), Inches(2.4), Pt(3), clr=color)
        # number
        txt(s, x + Inches(0.2), Inches(2.65), Inches(2), Inches(0.3),
            num, FONT_MONO, 11, color, bold=True)
        # phase name
        txt(s, x + Inches(0.2), Inches(2.95), Inches(2), Inches(0.45),
            name, FONT_HEADING, 17, NEUTRAL_950, bold=True)
        # description
        txt(s, x + Inches(0.2), Inches(3.55), Inches(2.0), Inches(0.8),
            desc, FONT_DISPLAY, 11, NEUTRAL_700, spacing=14)
        # divider
        bar(s, x + Inches(0.3), Inches(4.5), Inches(0.7), Pt(1), clr=NEUTRAL_200)
        # agents
        txt(s, x + Inches(0.2), Inches(4.6), Inches(2), Inches(0.25),
            "AGENTES", FONT_HEADING, 9, color, bold=True)
        txt(s, x + Inches(0.2), Inches(4.85), Inches(2), Inches(1.0),
            agents, FONT_MONO, 9, NEUTRAL_700, spacing=14)
        # rhythm
        txt(s, x + Inches(0.2), Inches(6.0), Inches(2), Inches(0.25),
            "RITMO", FONT_HEADING, 9, color, bold=True)
        txt(s, x + Inches(0.2), Inches(6.25), Inches(2), Inches(0.6),
            when, FONT_BODY, 10, NEUTRAL_700, spacing=14)
        # arrow
        if i < 4:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + Inches(2.4), Inches(4.45),
                                      Inches(0.18), Inches(0.22))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT_500
            arr.line.fill.background()
    footer(s, pg, TOTAL, dark=False)


# ---------------- 07-11 cada FASE em detalhe ----------------
def phase_slide(pg, num, name, color, dark, agents, when, prompt_template,
                expected_output, next_step, antipattern):
    s = prs.slides.add_slide(BLANK)
    if dark:
        set_bg(s, PRIMARY_950)
        chapter_num(s, num, clr=ACCENT_400)
        text_clr = WHITE
        body_clr = NEUTRAL_300
        sec_clr = ACCENT_500
        card_bg = PRIMARY_900
        small_label = ACCENT_400
    else:
        set_bg(s, WHITE)
        left_bar(s)
        section_label(s, num)
        text_clr = NEUTRAL_950
        body_clr = NEUTRAL_700
        sec_clr = ACCENT_600
        card_bg = NEUTRAL_50
        small_label = ACCENT_600

    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.6),
        f"FASE {num.split('.')[1]} — {name}",
        FONT_HEADING, 28, text_clr, bold=True)
    bar(s, Inches(1.2), Inches(1.75), Inches(1.0), Pt(3), clr=color)

    # Top row: agents + when
    rounded(s, Inches(1.2), Inches(2.1), Inches(5.6), Inches(1.0),
            card_bg, line=color if dark else NEUTRAL_200, line_w=Pt(0.5))
    txt(s, Inches(1.4), Inches(2.2), Inches(5.4), Inches(0.3),
        "AGENTES", FONT_HEADING, 10, sec_clr, bold=True)
    txt(s, Inches(1.4), Inches(2.5), Inches(5.4), Inches(0.6),
        agents, FONT_MONO, 13, text_clr, bold=True, spacing=18)

    rounded(s, Inches(7.1), Inches(2.1), Inches(5.6), Inches(1.0),
            card_bg, line=color if dark else NEUTRAL_200, line_w=Pt(0.5))
    txt(s, Inches(7.3), Inches(2.2), Inches(5.4), Inches(0.3),
        "QUANDO ACIONAR", FONT_HEADING, 10, sec_clr, bold=True)
    txt(s, Inches(7.3), Inches(2.5), Inches(5.4), Inches(0.6),
        when, FONT_BODY, 12, body_clr, spacing=16)

    # Prompt template
    txt(s, Inches(1.2), Inches(3.25), Inches(11), Inches(0.3),
        "PROMPT TEMPLATE  •  copie, ajuste e cole no Claude Code",
        FONT_HEADING, 10, small_label, bold=True)
    code_block(s, Inches(1.2), Inches(3.55), Inches(11.5), Inches(1.55),
                prompt_template, sz=11, clr=ACCENT_400, bg=PRIMARY_950,
                border=ACCENT_500)

    # Bottom row: expected output, next step, antipattern
    rounded(s, Inches(1.2), Inches(5.3), Inches(3.7), Inches(1.55),
            card_bg, line=SUCCESS, line_w=Pt(0.5))
    txt(s, Inches(1.4), Inches(5.4), Inches(3.5), Inches(0.25),
        "✓ SAIDA ESPERADA", FONT_HEADING, 9, SUCCESS, bold=True)
    txt(s, Inches(1.4), Inches(5.65), Inches(3.5), Inches(1.15),
        expected_output, FONT_BODY, 10, body_clr, spacing=14)

    rounded(s, Inches(5.1), Inches(5.3), Inches(3.7), Inches(1.55),
            card_bg, line=color, line_w=Pt(0.5))
    txt(s, Inches(5.3), Inches(5.4), Inches(3.5), Inches(0.25),
        "→ PROXIMO PASSO", FONT_HEADING, 9, sec_clr, bold=True)
    txt(s, Inches(5.3), Inches(5.65), Inches(3.5), Inches(1.15),
        next_step, FONT_BODY, 10, body_clr, spacing=14)

    rounded(s, Inches(9.0), Inches(5.3), Inches(3.7), Inches(1.55),
            card_bg, line=DANGER, line_w=Pt(0.5))
    txt(s, Inches(9.2), Inches(5.4), Inches(3.5), Inches(0.25),
        "✗ ANTI-PADRAO", FONT_HEADING, 9, DANGER, bold=True)
    txt(s, Inches(9.2), Inches(5.65), Inches(3.5), Inches(1.15),
        antipattern, FONT_BODY, 10, body_clr, spacing=14)

    footer(s, pg, TOTAL, dark=dark)


def s07_fase1_input(pg):
    phase_slide(
        pg=pg, num="03.1", name="INPUT", color=DANGER, dark=True,
        agents="academic-researcher\nsource-validator",
        when="2-3x por semana, ou\nquando o tema pede fontes\nnovas (gap, nova vertente).",
        prompt_template=(
            'Acione o academic-researcher para buscar fontes recentes\n'
            '(2020-2026) sobre {SEU TEMA} em SciELO, Scholar e CAPES.\n'
            'Foco em artigos com Qualis A1-B2. Retornar 8-12 fontes\n'
            'priorizadas com DOI verificavel.\n\n'
            'Em seguida, acione source-validator para auditar cada\n'
            'referencia retornada antes de eu salvar no biblio.bib.'
        ),
        expected_output=(
            "Lista em research/{data}/sources.md\n"
            "com 8-12 fontes (DOI, Qualis,\n"
            "resumo de 3 linhas, relevancia 1-5).\n"
            "Validacao com simbolos\n"
            "✓ / 🟡 / ❌ por afirmacao."
        ),
        next_step=(
            "1) Aprovar fontes ✓\n"
            "2) citation-manager: adicionar\n   .bib\n"
            "3) Mover PDFs aprovados\n   para Bibliografia/PDFs/\n"
            "4) knowledge-architect ficha"
        ),
        antipattern=(
            "Aceitar lista sem rodar\n"
            "source-validator. IA pode\n"
            "alucinar DOI, ano e ate\n"
            "autor. Toda fonte tem que\n"
            "passar pelo gatekeeper."
        ),
    )


def s08_fase2_desenho(pg):
    phase_slide(
        pg=pg, num="03.2", name="DESENHO", color=PRIMARY_700, dark=False,
        agents="methodology-advisor",
        when="No inicio do projeto.\nReavaliacao apos qualificacao\nou se mudar pergunta.",
        prompt_template=(
            'methodology-advisor: minha pergunta de pesquisa e\n'
            '"{PERGUNTA}". Pretendo investigar {OBJETO}. Tenho\n'
            '{TEMPO} ate qualificacao e {RESTRICAO_DE_CAMPO}.\n\n'
            'Apresente 3 alternativas metodologicas (qual/quant/misto)\n'
            'com paradigma, instrumentos, autores canonicos e os\n'
            '3 pontos mais provaveis de critica da banca.'
        ),
        expected_output=(
            "3 alternativas comparadas:\n"
            "paradigma, estrategia, amostra,\n"
            "instrumentos, autores de\n"
            "referencia. Recomendacao\n"
            "final justificada em 5 linhas."
        ),
        next_step=(
            "Voce escolhe uma.\n"
            "methodology-advisor preenche\n"
            "memory/project_methodology.md\n"
            "com paradigma + instrumentos +\n"
            "validade declarada."
        ),
        antipattern=(
            'Pedir "qual metodologia\n'
            'usar?" sem dar pergunta de\n'
            "pesquisa, restricoes e\n"
            "objetivos. O agente vai\n"
            "responder no vacuo."
        ),
    )


def s09_fase3_digestao(pg):
    phase_slide(
        pg=pg, num="03.3", name="DIGESTAO", color=ACCENT_600, dark=True,
        agents="knowledge-architect",
        when="A cada leitura concluida.\nRito quinzenal de promocao\nde notas raw → connected.",
        prompt_template=(
            'knowledge-architect: acabei de ler {AUTOR (ANO)}\n'
            '"{TITULO}". Aqui esta meu fichamento bruto:\n\n'
            '{COLE FICHAMENTO}\n\n'
            'Por favor:\n'
            '1) Crie as notas atomicas necessarias em\n'
            '   Notas-Atomicas/permanent/\n'
            '2) Linke bidirecionalmente com 2-5 notas existentes\n'
            '3) Atualize a MOC do tema se ja houver massa critica'
        ),
        expected_output=(
            "3-5 notas atomicas em\n"
            "Notas-Atomicas/permanent/\n"
            "{ts}-{slug}.md, cada uma com\n"
            "frontmatter YAML + 2 links\n"
            "bidirecionais minimos."
        ),
        next_step=(
            "Revisar links sugeridos\n"
            "(o agente as vezes erra).\n"
            "Promover status de raw para\n"
            "refined depois de 1 semana,\n"
            "para connected apos 2 links."
        ),
        antipattern=(
            "Aceitar nota com mais de\n"
            "uma ideia. Se tem dois\n"
            "argumentos independentes,\n"
            "PEDIR para dividir.\n"
            "Atomicidade e sagrada."
        ),
    )


def s10_fase4_producao(pg):
    phase_slide(
        pg=pg, num="03.4", name="PRODUCAO", color=PRIMARY_950, dark=False,
        agents="academic-writer\ncitation-manager",
        when="Sprints de redacao.\nQuando voce tiver 8-12 notas\nconectadas sobre um sub-tema.",
        prompt_template=(
            'academic-writer: redija o esqueleto da secao 2.3 do meu\n'
            'capitulo 2, partindo destas notas atomicas conectadas:\n'
            '[[202604201430-poder-disciplinar]]\n'
            '[[202604180900-zuboff-vigilancia]]\n'
            '[[...]]\n\n'
            'Use estrutura tese > evidencia > contraponto > sintese.\n'
            'Limite 1500 palavras. Marque [CITAR: ...] onde faltar fonte.\n'
            'Em seguida, citation-manager formata as referencias usadas.'
        ),
        expected_output=(
            "Rascunho .md em\n"
            "01_Projects/{atual}/drafts/\n"
            "com 1500 palavras, marcacoes\n"
            "[CITAR: ...] explicitas\n"
            "onde falta fonte."
        ),
        next_step=(
            "1) Revisar onde aparece\n"
            "   [CITAR: ...]\n"
            "2) Acionar academic-researcher\n"
            "   para preencher\n"
            "3) Iterar 1-2x antes de\n"
            "   acionar fase 5 (polimento)"
        ),
        antipattern=(
            "Pedir o capitulo inteiro\n"
            "de uma vez. Quebre por\n"
            "secao (1.500-2.000 palavras).\n"
            "Cada secao iterada, depois\n"
            "costure. Resultado e melhor."
        ),
    )


def s11_fase5_polimento(pg):
    phase_slide(
        pg=pg, num="03.5", name="POLIMENTO", color=SUCCESS, dark=True,
        agents="ai-pattern-auditor\nvoice-humanizer\npeer-reviewer\n(NESTA ORDEM)",
        when="ANTES de qualquer\nsubmissao: orientador,\nbanca, periodico.",
        prompt_template=(
            '# Sequencia obrigatoria — execute em ordem:\n\n'
            '1. source-validator: audite todas as citacoes do arquivo\n'
            '   drafts/cap2-secao3.md.\n\n'
            '2. ai-pattern-auditor: gere mapa de calor paragrafo a\n'
            '   paragrafo, comparando com memory/user_writing_voice.md\n\n'
            '3. voice-humanizer: reescreva os paragrafos 🔴 e 🟠\n'
            '   apontados pelo auditor, em colaboracao comigo.\n\n'
            '4. peer-reviewer: simule banca e gere 10 piores perguntas.'
        ),
        expected_output=(
            "Quatro outputs em sequencia:\n"
            "✓ validacao de citacoes\n"
            "✓ heatmap com score 0-10\n"
            "✓ draft reescrito na sua voz\n"
            "✓ relatorio peer + 10 perguntas"
        ),
        next_step=(
            "Submeter ao orientador.\n"
            "Salvar peer-review/{data}.md\n"
            "como historico. Atualizar\n"
            "memory/banca_questions.md\n"
            "para preparar a defesa."
        ),
        antipattern=(
            "Pular o auditor e ir direto\n"
            "ao humanizer. O auditor da\n"
            "FOCO para a reescrita —\n"
            "sem ele, voce reescreve\n"
            "no escuro e desperdica esforco."
        ),
    )


# ---------------- 12-15 RITMOS DE USO ----------------
def s12_ritmo_diario(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "04.1  RITMO DIARIO")
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "30 minutos por dia",
        FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.8), Inches(11), Inches(0.4),
        "Constancia importa mais que volume. Faca pouco mas todo dia.",
        FONT_BODY, 13, NEUTRAL_600)

    blocks = [
        ("0-5 min", "Capturar no Inbox",
         "Tudo que apareceu hoje: PDF, foto,\nsnippet, audio. Sem julgamento.",
         "00_Inbox/", DANGER),
        ("5-15 min", "Daily note",
         "Templater no Obsidian gera nota do dia\ncom data + ideias rapidas.",
         "Templater plugin", PRIMARY_700),
        ("15-25 min", "1 nota atomica nova",
         "knowledge-architect transforma a leitura\ndo dia em pelo menos 1 nota.",
         "knowledge-architect", ACCENT_600),
        ("25-30 min", "Revisar 1 nota antiga",
         "Abra uma nota com status connected\nde >30 dias. Atualize ou link com nova.",
         "Notas-Atomicas/permanent/", SUCCESS),
    ]
    for i, (time, name, desc, ref, color) in enumerate(blocks):
        x = Inches(0.7 + (i % 2) * 6.2)
        y = Inches(2.55 + (i // 2) * 2.15)
        rounded(s, x, y, Inches(6), Inches(2),
                NEUTRAL_50, line=NEUTRAL_200)
        bar(s, x, y, Inches(6), Pt(3), clr=color)
        # time badge
        rounded(s, x + Inches(0.3), y + Inches(0.3), Inches(1.4), Inches(0.4), color)
        txt(s, x + Inches(0.3), y + Inches(0.36), Inches(1.4), Inches(0.3),
            time, FONT_MONO, 11, WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(1.85), y + Inches(0.32), Inches(4), Inches(0.4),
            name, FONT_HEADING, 16, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.3), y + Inches(0.85), Inches(5.5), Inches(0.7),
            desc, FONT_BODY, 11, NEUTRAL_700, spacing=15)
        txt(s, x + Inches(0.3), y + Inches(1.55), Inches(5.5), Inches(0.3),
            "→ " + ref, FONT_MONO, 10, color, bold=True)

    # bottom CTA
    rounded(s, Inches(1.2), Inches(6.85), Inches(11.5), Inches(0.4), PRIMARY_950)
    txt(s, Inches(1.4), Inches(6.9), Inches(11.3), Inches(0.3),
        "TOTAL  •  30 minutos  •  Idealmente sempre no mesmo horario do dia",
        FONT_MONO, 11, ACCENT_400, bold=True)
    footer(s, pg, TOTAL, dark=False)


def s13_ritmo_semanal(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "04.2", clr=ACCENT_400)
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Ritmo SEMANAL — 1 hora aos domingos",
        FONT_HEADING, 28, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(1.75), Inches(1.2))
    txt(s, Inches(1.2), Inches(1.95), Inches(11), Inches(0.4),
        "Manutencao e consolidacao do que foi capturado durante a semana.",
        FONT_BODY, 13, NEUTRAL_400)

    blocks = [
        ("0-15", "Esvaziar Inbox",
         "Cada item: arquivar, fichar ou descartar.",
         "00_Inbox/ → 03_Resources/"),
        ("15-30", "Atualizar biblio.bib",
         "citation-manager processa novas fontes\ndo PDFs/ que ainda nao tem .bib.",
         "citation-manager"),
        ("30-45", "Promover notas",
         "raw (>7 dias) → refined.\nrefined com 2 links → connected.",
         "Notas-Atomicas/permanent/"),
        ("45-55", "Atualizar MOCs",
         "Se algum tema tem 8+ notas conectadas,\ncrie ou atualize Map of Content.",
         "Notas-Atomicas/index/"),
        ("55-60", "Planejar a semana",
         "Definir prioridade da semana:\nleitura, redacao, fonte, peer review.",
         "01_Projects/{atual}/README.md"),
    ]
    for i, (time, name, desc, ref) in enumerate(blocks):
        y = Inches(2.55 + i * 0.85)
        rounded(s, Inches(1.2), y, Inches(11.5), Inches(0.78),
                PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
        # time badge
        rounded(s, Inches(1.4), y + Inches(0.13), Inches(0.95), Inches(0.5), ACCENT_500)
        txt(s, Inches(1.4), y + Inches(0.22), Inches(0.95), Inches(0.3),
            time + "min", FONT_MONO, 10, PRIMARY_950, bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(2.55), y + Inches(0.08), Inches(3), Inches(0.35),
            name, FONT_HEADING, 14, WHITE, bold=True)
        txt(s, Inches(2.55), y + Inches(0.42), Inches(6), Inches(0.4),
            desc, FONT_BODY, 11, NEUTRAL_300, spacing=14)
        txt(s, Inches(8.7), y + Inches(0.25), Inches(3.9), Inches(0.4),
            "→ " + ref, FONT_MONO, 10, ACCENT_400)
    footer(s, pg, TOTAL)


def s14_ritmo_sprint(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "04.3  RITMO POR SPRINT")
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Sprint de capitulo — 1 a 2 semanas",
        FONT_HEADING, 28, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.8), Inches(11), Inches(0.4),
        "Quando voce decide redigir uma secao ou capitulo. Inicie quando ha massa critica de notas.",
        FONT_BODY, 13, NEUTRAL_600)

    days = [
        ("D1", "Selecionar notas",
         "Escolher 8-12 notas atomicas\nconectadas sobre o tema da secao.",
         "knowledge-architect"),
        ("D2", "Esqueleto",
         "academic-writer gera outline:\ntese de cada paragrafo + ordem.",
         "academic-writer"),
        ("D3-D5", "Redigir",
         "Iterar redacao secao a secao.\n1500-2000 palavras por dia.",
         "academic-writer"),
        ("D6", "Audit de fontes",
         "source-validator audita cada\ncitacao do draft.",
         "source-validator"),
        ("D7", "Mapa de IA",
         "ai-pattern-auditor gera heatmap.\nIdentificar paragrafos 🔴.",
         "ai-pattern-auditor"),
        ("D8-D10", "Humanizar",
         "voice-humanizer reescreve\nfocado nos 🔴 com sua voz.",
         "voice-humanizer"),
        ("D11", "Banca simulada",
         "peer-reviewer aponta lacunas\ne 10 piores perguntas.",
         "peer-reviewer"),
        ("D12", "Iterar correcoes",
         "Resolver criticas. Reaudit\ncom auditor + voice se preciso.",
         "ciclo polimento"),
        ("D14", "Submeter",
         "Versao final ao orientador\nou periodico.",
         "GO!"),
    ]
    for i, (d, name, desc, agent) in enumerate(days):
        col = i % 3
        row = i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.6 + row * 1.45)
        rounded(s, x, y, Inches(4), Inches(1.3),
                NEUTRAL_50, line=NEUTRAL_200)
        # day badge
        rounded(s, x, y, Inches(0.95), Inches(1.3), PRIMARY_950)
        txt(s, x, y + Inches(0.5), Inches(0.95), Inches(0.3),
            d, FONT_MONO, 13, ACCENT_400, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(1.05), y + Inches(0.1), Inches(2.85), Inches(0.3),
            name, FONT_HEADING, 12, NEUTRAL_950, bold=True)
        txt(s, x + Inches(1.05), y + Inches(0.4), Inches(2.85), Inches(0.6),
            desc, FONT_BODY, 9, NEUTRAL_700, spacing=13)
        txt(s, x + Inches(1.05), y + Inches(1.0), Inches(2.85), Inches(0.25),
            "→ " + agent, FONT_MONO, 8, ACCENT_600, bold=True)
    footer(s, pg, TOTAL, dark=False)


def s15_ritmo_submissao(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "04.4", clr=ACCENT_400)
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Pre-submissao — checklist final",
        FONT_HEADING, 28, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(1.75), Inches(1.2))
    txt(s, Inches(1.2), Inches(1.95), Inches(11), Inches(0.4),
        "Antes de enviar para banca, periodico ou orientador. NAO PULE NENHUMA.",
        FONT_BODY, 13, NEUTRAL_400)

    items = [
        ("01", "source-validator", "Toda citacao verificada (✓/🟡/❌). Zero alucinacoes."),
        ("02", "ai-pattern-auditor", "Score medio < 4.0. Zero paragrafos 🔴."),
        ("03", "voice-humanizer", "Cada paragrafo voce defende em 30s oralmente."),
        ("04", "peer-reviewer", "10 piores perguntas listadas e respondidas."),
        ("05", "citation-manager", "biblio.bib formatado na norma exigida."),
        ("06", "Leitura humana", "Voce releu integralmente. Soa como voce."),
        ("07", "Limites", "Palavras, paginas, formato dentro da especificacao."),
        ("08", "Etica", "CEP/CONEP aprovado se envolve humanos. TCLE incluido."),
        ("09", "Anti-plagio", "Submeter a Turnitin/CopySpider. Score < 15%."),
        ("10", "Backup", "Versao final salva em git e Drive antes de enviar."),
    ]
    for i, (num, name, desc) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.7 + col * 6.2)
        y = Inches(2.55 + row * 0.85)
        rounded(s, x, y, Inches(6), Inches(0.78),
                PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
        rounded(s, x, y, Inches(0.7), Inches(0.78), ACCENT_500)
        txt(s, x, y + Inches(0.22), Inches(0.7), Inches(0.4),
            num, FONT_MONO, 13, PRIMARY_950, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.85), y + Inches(0.07), Inches(5), Inches(0.35),
            name, FONT_HEADING, 12, WHITE, bold=True)
        txt(s, x + Inches(0.85), y + Inches(0.42), Inches(5), Inches(0.35),
            desc, FONT_BODY, 10, NEUTRAL_300)
    footer(s, pg, TOTAL)


# ---------------- 16 DECISION TREE ----------------
def s16_decision_tree(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "05  DECISION TREE")
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Em duvida? Use este fluxograma",
        FONT_HEADING, 28, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.8), Inches(11), Inches(0.4),
        "Pergunta inicial → resposta indica qual agente acionar agora.",
        FONT_BODY, 13, NEUTRAL_600)

    decisions = [
        ('"Tenho duvida sobre o\nmetodo da pesquisa?"',  "methodology-advisor", PRIMARY_700),
        ('"Preciso de fontes\nnovas sobre {tema}?"',     "academic-researcher", DANGER),
        ('"Recebi texto de IA —\nposso confiar?"',       "source-validator",    DANGER),
        ('"Acabei de ler um\nartigo importante."',       "knowledge-architect", ACCENT_600),
        ('"Tenho 8+ notas conec-\ntadas, posso escrever?"', "academic-writer",  PRIMARY_950),
        ('"Como cito esta obra\nem ABNT?"',              "citation-manager",   NEUTRAL_700),
        ('"Texto pronto. Onde\nestao as marcas de IA?"', "ai-pattern-auditor", DANGER),
        ('"Quero reescrever\nna minha voz."',            "voice-humanizer",    ACCENT_600),
        ('"Vou submeter — qual\na critica mais dura?"',  "peer-reviewer",      SUCCESS),
    ]
    for i, (question, agent, color) in enumerate(decisions):
        col = i % 3
        row = i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.55 + row * 1.55)
        rounded(s, x, y, Inches(4), Inches(1.4),
                NEUTRAL_50, line=NEUTRAL_200)
        bar(s, x, y, Inches(4), Pt(3), clr=color)
        # question
        txt(s, x + Inches(0.2), y + Inches(0.2), Inches(3.7), Inches(0.65),
            question, FONT_DISPLAY, 12, NEUTRAL_950, spacing=15)
        # arrow
        arr = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                  x + Inches(1.85), y + Inches(0.85),
                                  Inches(0.3), Inches(0.18))
        arr.fill.solid()
        arr.fill.fore_color.rgb = color
        arr.line.fill.background()
        # agent
        txt(s, x + Inches(0.2), y + Inches(1.05), Inches(3.7), Inches(0.3),
            agent, FONT_MONO, 11, color, bold=True, align=PP_ALIGN.CENTER)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 17 CHEAT SHEET ----------------
def s17_cheatsheet(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "06", clr=ACCENT_400)
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Cheat sheet — todos os 9 agentes",
        FONT_HEADING, 28, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(1.75), Inches(1.2))
    txt(s, Inches(1.2), Inches(1.95), Inches(11), Inches(0.4),
        "Imprimir e deixar perto do computador.",
        FONT_BODY, 13, NEUTRAL_400)

    headers = ["AGENTE", "QUANDO ACIONAR", "PROMPT-CHAVE"]
    cols = [Inches(2.4), Inches(3.3), Inches(6.8)]
    cx = [Inches(0.7), Inches(3.1), Inches(6.4)]
    y_h = Inches(2.55)
    filled_rect(s, Inches(0.7), y_h, Inches(12.5), Inches(0.4), ACCENT_500)
    for i, h in enumerate(headers):
        txt(s, cx[i] + Inches(0.1), y_h + Inches(0.07), cols[i], Inches(0.3),
            h, FONT_HEADING, 10, PRIMARY_950, bold=True)

    rows = [
        ("methodology-advisor",   "Inicio + qualif.",     "metodologia para {pergunta}, 3 alternativas"),
        ("academic-researcher",   "Fontes novas",         "buscar {tema} em SciELO/Scholar 2020-2026"),
        ("source-validator",      "Pos qualquer IA",      "audite as citacoes deste arquivo"),
        ("knowledge-architect",   "Pos cada leitura",     "fichar {autor (ano)} e gerar notas atomicas"),
        ("academic-writer",       "Sprint redacao",       "redija secao X a partir destas notas: [[...]]"),
        ("citation-manager",      "Toda nova fonte",      "adicione {DOI} ao biblio.bib em ABNT"),
        ("ai-pattern-auditor",    "Antes do humanizer",   "mapa de calor de marcas de IA no draft"),
        ("voice-humanizer",       "1a sessao + pre-sub.", "aplicar questionario / humanizar 🔴 do auditor"),
        ("peer-reviewer",         "Antes da banca",       "simule banca, 10 piores perguntas"),
    ]
    for r, (a, w, p) in enumerate(rows):
        y = Inches(3.0 + r * 0.42)
        if r % 2 == 0:
            filled_rect(s, Inches(0.7), y, Inches(12.5), Inches(0.42), PRIMARY_900)
        txt(s, cx[0] + Inches(0.1), y + Inches(0.1), cols[0], Inches(0.3),
            a, FONT_MONO, 11, ACCENT_400, bold=True)
        txt(s, cx[1] + Inches(0.1), y + Inches(0.1), cols[1], Inches(0.3),
            w, FONT_BODY, 10, NEUTRAL_300)
        txt(s, cx[2] + Inches(0.1), y + Inches(0.1), cols[2], Inches(0.3),
            p, FONT_MONO, 9, WHITE)
    footer(s, pg, TOTAL)


# ---------------- 18 CLOSING ----------------
def s18_closing(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    filled_rect(s, Inches(0), Inches(2.6), SLIDE_W, Inches(0.04), ACCENT_500)

    txt(s, Inches(1.2), Inches(1.0), Inches(11), Inches(0.4),
        "PROXIMOS 3 PASSOS HOJE", FONT_HEADING, 12, ACCENT_400, bold=True)
    txt(s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.9),
        "Comece pelo fundamento.", FONT_DISPLAY, 38, WHITE, bold=True)
    txt(s, Inches(1.2), Inches(2.15), Inches(11), Inches(0.35),
        "Tres acoes em ordem para destravar todo o resto.",
        FONT_DISPLAY, 16, ACCENT_400)

    steps = [
        ("01", "Aplique o questionario do voice-humanizer",
         "30 minutos. Sem isso, todo polimento opera no escuro.\nCole no Claude Code: 'Acione voice-humanizer para aplicar o questionario inicial'"),
        ("02", "Defina sua pergunta de pesquisa em memory/project_thesis.md",
         "Uma frase clara. Se ainda nao tem, acione methodology-advisor para te ajudar a formular."),
        ("03", "Capture suas primeiras 5 fontes",
         "academic-researcher busca, source-validator audita, knowledge-architect ficha.\nTeste o ciclo end-to-end com fontes reais do seu tema."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(3.0 + i * 1.0)
        rounded(s, Inches(1.2), y, Inches(0.7), Inches(0.7), ACCENT_500)
        txt(s, Inches(1.2), y + Inches(0.18), Inches(0.7), Inches(0.4),
            num, FONT_MONO, 16, PRIMARY_950, bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(2.1), y + Inches(0.1), Inches(11), Inches(0.4),
            title, FONT_HEADING, 16, WHITE, bold=True)
        txt(s, Inches(2.1), y + Inches(0.45), Inches(10.5), Inches(0.5),
            desc, FONT_BODY, 11, NEUTRAL_300, spacing=14)

    txt(s, Inches(1.2), Inches(6.85), Inches(11), Inches(0.3),
        "Manual completo: github.com/byfabioviana/agent-study-assistant",
        FONT_MONO, 10, ACCENT_400)
    txt(s, Inches(1.2), Inches(7.15), Inches(11), Inches(0.3),
        "Verticalis AI Journey  •  Education  •  Maio 2026  •  v1.0",
        FONT_MONO, 9, NEUTRAL_500, align=PP_ALIGN.RIGHT)


# ============================================================================
# build all
# ============================================================================
s01_cover()
s02_toc(2)
s03_manifesto(3)
s04_prerequisito(4)
s05_anatomia(5)
s06_framework(6)
s07_fase1_input(7)
s08_fase2_desenho(8)
s09_fase3_digestao(9)
s10_fase4_producao(10)
s11_fase5_polimento(11)
s12_ritmo_diario(12)
s13_ritmo_semanal(13)
s14_ritmo_sprint(14)
s15_ritmo_submissao(15)
s16_decision_tree(16)
s17_cheatsheet(17)
s18_closing(18)

out = ROOT / "docs" / "manual-uso-agentes.pptx"
prs.save(str(out))
print(f"PPTX gerado: {out}")
print(f"Total de slides: {len(prs.slides)}")
