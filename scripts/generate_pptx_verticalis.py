"""
Sistema de Curadoria Academica - Verticalis AI Journey deck
Padrao Oxford Heritage v2.0 | Plus Jakarta Sans + Lora + IBM Plex Mono
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# === Oxford Heritage v2.0 tokens ===
PRIMARY_950 = RGBColor(0x0B, 0x28, 0x23)
PRIMARY_900 = RGBColor(0x2D, 0x45, 0x3E)
PRIMARY_700 = RGBColor(0x3A, 0x56, 0x4E)
PRIMARY_600 = RGBColor(0x2D, 0x6E, 0x43)
ACCENT_600 = RGBColor(0xC5, 0xB3, 0x58)
ACCENT_500 = RGBColor(0xC5, 0xB3, 0x58)
ACCENT_400 = RGBColor(0xD4, 0xC8, 0x78)
NEUTRAL_950 = RGBColor(0x1E, 0x1E, 0x1E)
NEUTRAL_800 = RGBColor(0x33, 0x33, 0x33)
NEUTRAL_700 = RGBColor(0x4A, 0x4A, 0x4A)
NEUTRAL_600 = RGBColor(0x66, 0x66, 0x66)
NEUTRAL_500 = RGBColor(0x88, 0x88, 0x88)
NEUTRAL_400 = RGBColor(0xAA, 0xAA, 0xAA)
NEUTRAL_300 = RGBColor(0xCC, 0xCC, 0xCC)
NEUTRAL_200 = RGBColor(0xE0, 0xDD, 0xD5)
NEUTRAL_50 = RGBColor(0xF4, 0xF1, 0xEA)
DEW_GREEN = RGBColor(0xF4, 0xF1, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DANGER = RGBColor(0xC0, 0x39, 0x2B)
SUCCESS = RGBColor(0x2D, 0x6E, 0x43)
WARN = RGBColor(0xD4, 0x8B, 0x1F)

FONT_HEADING = "Plus Jakarta Sans"
FONT_BODY = "Plus Jakarta Sans"
FONT_DISPLAY = "Lora"
FONT_MONO = "IBM Plex Mono"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

ROOT = Path(__file__).resolve().parent.parent
LOGO_DIR = Path("D:/corporateProjects/verticalis-ai-journey/assets/brand/logo")
LOGO_HORIZ_LIGHT = LOGO_DIR / "horizontal" / "logo-horizontal-light.png"
LOGO_HORIZ_DARK = LOGO_DIR / "horizontal" / "logo-horizontal-dark.png"
MARK_LIGHT = LOGO_DIR / "mark" / "mark-light-512.png"
MARK_DARK = LOGO_DIR / "mark" / "mark-dark-512.png"

# ----------------------- helpers -----------------------
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def txt(slide, l, t, w, h, text, font=FONT_BODY, sz=18, clr=WHITE,
        bold=False, align=PP_ALIGN.LEFT, spacing=None):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.alignment = align
    if spacing:
        p.line_spacing = Pt(spacing)
    return box, tf


def bar(slide, l, t, w=Inches(1), h=Pt(3), clr=ACCENT_500):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = clr
    s.line.fill.background()
    return s


def left_bar(slide, clr=ACCENT_500):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(6), SLIDE_H)
    s.fill.solid()
    s.fill.fore_color.rgb = clr
    s.line.fill.background()


def section_label(slide, text, l=Inches(1.2), t=Inches(0.6), clr=ACCENT_600):
    txt(slide, l, t, Inches(8), Inches(0.35), text.upper(),
        FONT_HEADING, 11, clr, bold=True)


def chapter_num(slide, num, clr=ACCENT_400):
    txt(slide, Inches(1.2), Inches(0.6), Inches(1), Inches(0.35),
        num, FONT_MONO, 13, clr, bold=True)


def footer(slide, pg, total, dark=True):
    c = NEUTRAL_500 if dark else NEUTRAL_600
    txt(slide, Inches(1.2), Inches(7.05), Inches(7), Inches(0.35),
        "Sistema de Curadoria Academica  |  Verticalis AI Journey  |  Confidencial",
        FONT_BODY, 9, c)
    txt(slide, Inches(11.0), Inches(7.05), Inches(1.5), Inches(0.35),
        f"{pg:02d} / {total:02d}", FONT_MONO, 9, c, align=PP_ALIGN.RIGHT)


def add_logo(slide, kind="horiz_light", x=Inches(11.0), y=Inches(0.5), h=Inches(0.4)):
    """Logomarcas desativadas por design — apenas paleta/fontes/padrao do brand book."""
    return


def filled_rect(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = line_w
    return s


def rounded(slide, l, t, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = line_w
    return s


def icon_circle(slide, cx, cy, r, fill, mark, mark_color, mark_size=18, mark_font=FONT_MONO):
    """Circulo com simbolo dentro como icone."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - r, cy - r, r*2, r*2)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    txt(slide, cx - r, cy - r + Emu(int(r*0.15)), r*2, r*2,
        mark, mark_font, mark_size, mark_color, bold=True, align=PP_ALIGN.CENTER)


# ============================================================================
# BUILD
# ============================================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
TOTAL = 22
BLANK = prs.slide_layouts[6]


# ---------------- 01. CAPA ----------------
def s01_cover():
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    # Bottom-right accent block
    filled_rect(s, Inches(8.5), Inches(4.5), Inches(4.8), Inches(3), PRIMARY_900)
    # Logo no canto superior direito
    add_logo(s, "horiz_light", Inches(10.5), Inches(0.5), Inches(0.45))

    txt(s, Inches(1.2), Inches(1.6), Inches(10), Inches(0.5),
        "VERTICALIS AI JOURNEY  /  EDUCATION", FONT_HEADING, 12, ACCENT_400, bold=True)
    bar(s, Inches(1.2), Inches(2.2), Inches(1.2), Pt(3))
    txt(s, Inches(1.2), Inches(2.5), Inches(11), Inches(1.5),
        "Sistema de Curadoria\nAcademica", FONT_DISPLAY, 52, WHITE, bold=True,
        spacing=58)
    txt(s, Inches(1.2), Inches(4.7), Inches(10), Inches(0.5),
        "PARA + Zettelkasten + CODE  |  Google Drive  |  Squad de 9 Agentes",
        FONT_BODY, 16, NEUTRAL_300)
    txt(s, Inches(1.2), Inches(5.3), Inches(10), Inches(0.4),
        "Pesquisa rapida, contexto persistente, autoria autentica.",
        FONT_DISPLAY, 17, ACCENT_400)
    txt(s, Inches(1.2), Inches(6.5), Inches(10), Inches(0.4),
        "Maio 2026  •  Confidencial  •  v1.0",
        FONT_MONO, 11, NEUTRAL_500)


# ---------------- 02. INDICE ----------------
def s02_toc(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "INDICE")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Conteudo", FONT_HEADING, 38, NEUTRAL_950, bold=True)
    chapters = [
        ("01", "Diagnostico", "As 3 friccoes que estamos resolvendo", False),
        ("02", "Framework", "PARA + Zettelkasten + CODE — overview", False),
        ("",   "  • 02.1 PARA", "Organizacao por acionabilidade", True),
        ("",   "  • 02.2 Zettelkasten", "Notas atomicas linkadas", True),
        ("",   "  • 02.3 CODE", "Workflow de processamento", True),
        ("03", "Plataforma", "Google Drive como infraestrutura", False),
        ("04", "Estrutura de pastas", "Arvore PARA numerada 00-99", False),
        ("05", "Nomenclatura", "Convencao YYYYMMDD-tipo-slug", False),
        ("06", "Workflow CODE", "Capture > Organize > Distill > Express", False),
        ("07", "Squad de 9 agentes", "Como cada um atua no Drive", False),
        ("08", "Voice + Auditor de IA", "Calibracao e diagnostico forense", False),
        ("09", "Pesquisa rapida", "5 estrategias de busca no vault", False),
        ("10", "Interface Obsidian", "Tier 1 — UI moderna sobre o Drive", False),
        ("11", "Adocao", "Cronograma 4 semanas e proximos passos", False),
    ]
    for i, (num, title, desc, sub) in enumerate(chapters):
        y = Inches(2.0 + i * 0.42)
        if num:
            txt(s, Inches(1.2), y, Inches(0.6), Inches(0.4),
                num, FONT_MONO, 13, ACCENT_500, bold=True)
        title_color = NEUTRAL_700 if sub else NEUTRAL_950
        title_size = 13 if sub else 15
        title_bold = False if sub else True
        title_font = FONT_BODY if sub else FONT_HEADING
        txt(s, Inches(2.0), y, Inches(4.2), Inches(0.4),
            title, title_font, title_size, title_color, bold=title_bold)
        desc_color = NEUTRAL_500 if sub else NEUTRAL_600
        desc_size = 11 if sub else 13
        txt(s, Inches(6.2), y, Inches(7), Inches(0.4),
            desc, FONT_BODY, desc_size, desc_color)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 03. MANIFESTO / PROBLEMA ----------------
def s03_manifesto(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "01")
    txt(s, Inches(1.2), Inches(1.0), Inches(11), Inches(1.0),
        '"O conhecimento que voce nao consegue\nrecuperar, voce nao tem."',
        FONT_DISPLAY, 36, WHITE, bold=True, spacing=44)
    bar(s, Inches(1.2), Inches(2.7), Inches(1.2))

    left = (
        "Sobram fontes. Sobram conversas com IA.\n"
        "Sobram PDFs baixados que nunca sao lidos.\n"
        "O que falta e o sistema que cura, valida\n"
        "e devolve esse material quando voce precisa.\n\n"
        "O mestrando trata curadoria como \"depois eu organizo\".\n"
        "Tratamos como infraestrutura."
    )
    right = (
        "Nao prometemos produtividade. Prometemos rigor.\n"
        "Nao prometemos atalho. Prometemos metodo.\n"
        "Nao prometemos enganar detector de IA.\n"
        "Prometemos texto que reflete sua voz — porque\n"
        "voce intervem em cada paragrafo, e existe um\n"
        "agente especialista em cada etapa do processo."
    )
    txt(s, Inches(1.2), Inches(3.2), Inches(5.5), Inches(3.2),
        left, FONT_BODY, 14, NEUTRAL_300, spacing=22)
    txt(s, Inches(7.0), Inches(3.2), Inches(5.5), Inches(3.2),
        right, FONT_BODY, 14, NEUTRAL_300, spacing=22)

    txt(s, Inches(1.2), Inches(6.5), Inches(11), Inches(0.4),
        "Inteligencia que se constroi — aplicada a sua jornada de mestrado.",
        FONT_DISPLAY, 16, ACCENT_400, align=PP_ALIGN.CENTER)
    footer(s, pg, TOTAL)


# ---------------- 04. DIAGNOSTICO (3 cards) ----------------
def s04_diagnostico(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "01  DIAGNOSTICO")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Tres friccoes que travam a jornada", FONT_HEADING, 32, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Identificadas em sessoes reais com mestrandos usando IA como copilot de pesquisa.",
        FONT_BODY, 14, NEUTRAL_600)

    cards = [
        ("01", "Curadoria pulverizada",
         "Resultados misturam fontes confiaveis,\nblogs e alucinacoes.",
         "Falta gatekeeper que valide cada\nafirmacao contra fonte primaria."),
        ("02", "Conhecimento volatil",
         "Insights ficam em conversas dispersas,\nsem estrutura cumulativa.",
         "Falta camada Zettelkasten / segundo\ncerebro persistente."),
        ("03", "Texto soa de IA",
         "Detectores acusam padrao estatistico.\nBaixa burstiness, vocabulario previsivel.",
         "Texto gerado em uma unica passada\nsem reescrita autoral ancorada na voz."),
    ]
    for i, (num, title, sintoma, causa) in enumerate(cards):
        x = Inches(1.2 + i * 4.0)
        rounded(s, x, Inches(2.5), Inches(3.7), Inches(4.3), DEW_GREEN,
                line=NEUTRAL_200, line_w=Pt(1))
        # Number bar
        bar(s, x + Inches(0.3), Inches(2.8), Inches(0.6), Pt(3))
        txt(s, x + Inches(0.3), Inches(2.95), Inches(2), Inches(0.3),
            num, FONT_MONO, 11, ACCENT_600, bold=True)
        txt(s, x + Inches(0.3), Inches(3.3), Inches(3.2), Inches(0.5),
            title, FONT_HEADING, 18, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.3), Inches(4.1), Inches(3.2), Inches(0.3),
            "SINTOMA", FONT_HEADING, 9, DANGER, bold=True)
        txt(s, x + Inches(0.3), Inches(4.4), Inches(3.2), Inches(1.0),
            sintoma, FONT_BODY, 12, NEUTRAL_700, spacing=18)
        txt(s, x + Inches(0.3), Inches(5.5), Inches(3.2), Inches(0.3),
            "CAUSA RAIZ", FONT_HEADING, 9, ACCENT_600, bold=True)
        txt(s, x + Inches(0.3), Inches(5.8), Inches(3.2), Inches(1.0),
            causa, FONT_BODY, 12, NEUTRAL_700, spacing=18)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 05. FRAMEWORK ----------------
def s05_framework(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "02")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Tres frameworks combinados", FONT_HEADING, 32, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "Cada framework cobre uma camada do problema. Juntos formam a infraestrutura.",
        FONT_BODY, 14, NEUTRAL_400)

    boxes = [
        ("PARA", "Tiago Forte", "Organizacao por nivel\nde acionabilidade",
         ["Projects: tem prazo", "Areas: continuas", "Resources: tematicos", "Archive: concluido"]),
        ("ZETTELKASTEN", "Niklas Luhmann", "Notas atomicas\nlinkadas",
         ["Uma ideia por nota", "Voz do aluno", "Links bidirecionais", "Crescimento organico"]),
        ("CODE", "Tiago Forte", "Workflow de\nprocessamento",
         ["Capture (inbox)", "Organize (PARA)", "Distill (Zettel)", "Express (drafts)"]),
    ]
    for i, (name, author, desc, bullets) in enumerate(boxes):
        x = Inches(1.2 + i * 4.0)
        rounded(s, x, Inches(2.9), Inches(3.7), Inches(3.9), PRIMARY_900,
                line=ACCENT_500, line_w=Pt(0.75))
        bar(s, x + Inches(0.3), Inches(3.15), Inches(0.6))
        txt(s, x + Inches(0.3), Inches(3.3), Inches(3.2), Inches(0.5),
            name, FONT_HEADING, 22, ACCENT_500, bold=True)
        txt(s, x + Inches(0.3), Inches(3.85), Inches(3.2), Inches(0.3),
            author, FONT_MONO, 10, NEUTRAL_400)
        txt(s, x + Inches(0.3), Inches(4.25), Inches(3.2), Inches(0.7),
            desc, FONT_DISPLAY, 14, WHITE, spacing=20)
        for j, b in enumerate(bullets):
            yj = Inches(5.05 + j * 0.36)
            # bullet dot
            d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.35), yj + Inches(0.08),
                                    Inches(0.08), Inches(0.08))
            d.fill.solid()
            d.fill.fore_color.rgb = ACCENT_500
            d.line.fill.background()
            txt(s, x + Inches(0.55), yj, Inches(3.0), Inches(0.3),
                b, FONT_BODY, 12, NEUTRAL_300)
    footer(s, pg, TOTAL)


# ---------------- 05.A PARA detalhado ----------------
def s05a_para(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "02.1  FRAMEWORK 1 DE 3")
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "PARA — Organizacao por acionabilidade",
        FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.75), Inches(11), Inches(0.4),
        "Tiago Forte  •  Building a Second Brain (2017)  •  Sistema universal de gestao de informacao digital",
        FONT_BODY, 13, NEUTRAL_600)

    # Hero: principio central
    rounded(s, Inches(1.2), Inches(2.35), Inches(11.5), Inches(0.85),
            PRIMARY_950)
    txt(s, Inches(1.5), Inches(2.5), Inches(11), Inches(0.3),
        "PRINCIPIO CENTRAL", FONT_HEADING, 10, ACCENT_400, bold=True)
    txt(s, Inches(1.5), Inches(2.78), Inches(11), Inches(0.4),
        "Organize informacao por acionabilidade — quao proximo esta da execucao — nao por origem ou tipo.",
        FONT_DISPLAY, 14, WHITE, bold=True)

    # 4 niveis em cards
    levels = [
        ("P", "PROJECTS",   "Tem prazo, fim claro",
         "Capitulo 2 da dissertacao\nArtigo para periodico\nPreparacao da qualificacao",
         DANGER),
        ("A", "AREAS",      "Continuas, sem fim",
         "Sua linha de pesquisa\nRelacao com orientador\nDisciplinas em curso",
         ACCENT_600),
        ("R", "RESOURCES",  "Tematicos, consultaveis",
         "Bibliografia do campo\nTemplates de fichamento\nGlossario de conceitos",
         PRIMARY_700),
        ("A", "ARCHIVE",    "Inativos ou concluidos",
         "Trabalhos de TCC\nDisciplinas finalizadas\nIdeias engavetadas",
         NEUTRAL_600),
    ]
    for i, (letter, name, sub, examples, color) in enumerate(levels):
        x = Inches(1.2 + i * 2.95)
        rounded(s, x, Inches(3.4), Inches(2.8), Inches(2.65),
                NEUTRAL_50, line=NEUTRAL_200)
        bar(s, x, Inches(3.4), Inches(2.8), Pt(3), clr=color)
        # letter badge
        rounded(s, x + Inches(0.2), Inches(3.55), Inches(0.45), Inches(0.45), color)
        txt(s, x + Inches(0.2), Inches(3.62), Inches(0.45), Inches(0.4),
            letter, FONT_HEADING, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.75), Inches(3.55), Inches(2), Inches(0.3),
            name, FONT_HEADING, 13, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.75), Inches(3.85), Inches(2), Inches(0.25),
            sub, FONT_MONO, 9, color, bold=True)
        txt(s, x + Inches(0.2), Inches(4.3), Inches(2.6), Inches(1.6),
            examples, FONT_BODY, 11, NEUTRAL_700, spacing=15)

    # Bottom: quando funciona / quando nao + recomendacao
    rounded(s, Inches(1.2), Inches(6.2), Inches(3.7), Inches(0.7),
            DEW_GREEN, line=SUCCESS)
    txt(s, Inches(1.4), Inches(6.27), Inches(3.5), Inches(0.25),
        "QUANDO FUNCIONA", FONT_HEADING, 9, SUCCESS, bold=True)
    txt(s, Inches(1.4), Inches(6.5), Inches(3.5), Inches(0.35),
        "Multiplos projetos paralelos com prazos.",
        FONT_BODY, 10, NEUTRAL_700)

    rounded(s, Inches(5.1), Inches(6.2), Inches(3.7), Inches(0.7),
            DEW_GREEN, line=DANGER)
    txt(s, Inches(5.3), Inches(6.27), Inches(3.5), Inches(0.25),
        "QUANDO E LIMITADO", FONT_HEADING, 9, DANGER, bold=True)
    txt(s, Inches(5.3), Inches(6.5), Inches(3.5), Inches(0.35),
        "Nao captura conexao entre ideias — so onde ficam.",
        FONT_BODY, 10, NEUTRAL_700)

    rounded(s, Inches(9.0), Inches(6.2), Inches(3.7), Inches(0.7),
            PRIMARY_950)
    txt(s, Inches(9.2), Inches(6.27), Inches(3.5), Inches(0.25),
        "PARA O MESTRANDO", FONT_HEADING, 9, ACCENT_400, bold=True)
    txt(s, Inches(9.2), Inches(6.5), Inches(3.5), Inches(0.35),
        "Camada macro — onde as coisas vivem no Drive.",
        FONT_BODY, 10, NEUTRAL_300)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 05.B Zettelkasten detalhado ----------------
def s05b_zettel(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "02.2", clr=ACCENT_400)
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Zettelkasten — Notas atomicas linkadas",
        FONT_HEADING, 30, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(1.85), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.0), Inches(11), Inches(0.4),
        "Niklas Luhmann  •  sociologo alemao (1927-1998)  •  ~90.000 fichas, 70+ livros e 400+ artigos em 30 anos",
        FONT_BODY, 13, NEUTRAL_400)

    # Hero: principio central
    rounded(s, Inches(1.2), Inches(2.6), Inches(11.5), Inches(0.85),
            PRIMARY_900, line=ACCENT_500)
    txt(s, Inches(1.5), Inches(2.75), Inches(11), Inches(0.3),
        "PRINCIPIO CENTRAL", FONT_HEADING, 10, ACCENT_500, bold=True)
    txt(s, Inches(1.5), Inches(3.03), Inches(11), Inches(0.4),
        "Conhecimento gera valor pela conexao entre ideias atomicas — nao pelo acumulo.",
        FONT_DISPLAY, 14, WHITE, bold=True)

    # Esquerda: 5 principios
    txt(s, Inches(1.2), Inches(3.7), Inches(5.5), Inches(0.3),
        "OS 5 PRINCIPIOS DE LUHMANN", FONT_HEADING, 11, ACCENT_500, bold=True)
    principles = [
        ("01", "Atomicidade",   "Uma ideia por nota. Se tem dois argumentos independentes, divida."),
        ("02", "Autonomia",     "A nota deve ser inteligivel sem precisar abrir a fonte original."),
        ("03", "Conectividade", "Toda nota tem ao menos um link bidirecional para outra."),
        ("04", "Voz propria",   "Reescrita autoral — parafrasear e processar, nao copiar."),
        ("05", "Crescimento organico", "Sem hierarquia rigida — a estrutura emerge das conexoes."),
    ]
    for i, (num, name, desc) in enumerate(principles):
        y = Inches(4.0 + i * 0.55)
        txt(s, Inches(1.2), y, Inches(0.5), Inches(0.3),
            num, FONT_MONO, 11, ACCENT_400, bold=True)
        txt(s, Inches(1.7), y, Inches(2), Inches(0.3),
            name, FONT_HEADING, 12, WHITE, bold=True)
        txt(s, Inches(3.5), y + Inches(0.02), Inches(3.3), Inches(0.5),
            desc, FONT_BODY, 10, NEUTRAL_300, spacing=14)

    # Direita: tipos de nota + diferencial
    rounded(s, Inches(7.1), Inches(3.7), Inches(5.5), Inches(2.3),
            PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
    txt(s, Inches(7.3), Inches(3.85), Inches(5), Inches(0.3),
        "TIPOS DE NOTA", FONT_HEADING, 11, ACCENT_500, bold=True)
    types = [
        ("Literatura", "Fichamento por fonte (citekey).md"),
        ("Permanente", "Nucleo — uma ideia atomica autoral"),
        ("MOC",        "Map of Content — indice tematico curado"),
    ]
    for i, (n, d) in enumerate(types):
        y = Inches(4.2 + i * 0.55)
        rounded(s, Inches(7.3), y, Inches(0.18), Inches(0.4), ACCENT_500)
        txt(s, Inches(7.6), y, Inches(2), Inches(0.3),
            n, FONT_HEADING, 12, WHITE, bold=True)
        txt(s, Inches(7.6), y + Inches(0.3), Inches(4.8), Inches(0.3),
            d, FONT_BODY, 10, NEUTRAL_400)

    # Bottom: quando funciona / quando nao + recomendacao
    rounded(s, Inches(1.2), Inches(6.2), Inches(3.7), Inches(0.7),
            PRIMARY_900, line=SUCCESS)
    txt(s, Inches(1.4), Inches(6.27), Inches(3.5), Inches(0.25),
        "QUANDO FUNCIONA", FONT_HEADING, 9, SUCCESS, bold=True)
    txt(s, Inches(1.4), Inches(6.5), Inches(3.5), Inches(0.35),
        "Para gerar pensamento original em pesquisa.",
        FONT_BODY, 10, NEUTRAL_300)

    rounded(s, Inches(5.1), Inches(6.2), Inches(3.7), Inches(0.7),
            PRIMARY_900, line=DANGER)
    txt(s, Inches(5.3), Inches(6.27), Inches(3.5), Inches(0.25),
        "QUANDO E LIMITADO", FONT_HEADING, 9, DANGER, bold=True)
    txt(s, Inches(5.3), Inches(6.5), Inches(3.5), Inches(0.35),
        "Curva de aprendizado alta. Exige disciplina diaria.",
        FONT_BODY, 10, NEUTRAL_300)

    rounded(s, Inches(9.0), Inches(6.2), Inches(3.7), Inches(0.7),
            ACCENT_500)
    txt(s, Inches(9.2), Inches(6.27), Inches(3.5), Inches(0.25),
        "PARA O MESTRANDO", FONT_HEADING, 9, PRIMARY_950, bold=True)
    txt(s, Inches(9.2), Inches(6.5), Inches(3.5), Inches(0.35),
        "Nucleo da dissertacao — onde a tese se forma.",
        FONT_BODY, 10, PRIMARY_950)
    footer(s, pg, TOTAL)


# ---------------- 05.C CODE detalhado ----------------
def s05c_code_detailed(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "02.3  FRAMEWORK 3 DE 3")
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "CODE — Workflow de processamento",
        FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.75), Inches(11), Inches(0.4),
        "Tiago Forte  •  Building a Second Brain (2022)  •  Fluxo de informacao em 4 etapas",
        FONT_BODY, 13, NEUTRAL_600)

    # Hero
    rounded(s, Inches(1.2), Inches(2.35), Inches(11.5), Inches(0.85),
            PRIMARY_950)
    txt(s, Inches(1.5), Inches(2.5), Inches(11), Inches(0.3),
        "PRINCIPIO CENTRAL", FONT_HEADING, 10, ACCENT_400, bold=True)
    txt(s, Inches(1.5), Inches(2.78), Inches(11), Inches(0.4),
        "Informacao so vale se vira output. Cada etapa transforma input bruto em produto autoral.",
        FONT_DISPLAY, 14, WHITE, bold=True)

    # 4 etapas em cards verticais com profundidade
    stages = [
        ("C", "CAPTURE", "Salvar o que ressoa",
         "Drive Inbox\nclip de Scholar\naudio dictado",
         "Volume baixo curado > volume alto descartado.",
         DANGER),
        ("O", "ORGANIZE", "Armazenar por acionabilidade",
         "PARA aplica aqui:\nProjects/Areas/\nResources/Archive",
         "Pasta certa = recuperacao em segundos.",
         ACCENT_600),
        ("D", "DISTILL", "Condensar ao essencial",
         "Notas progressivas\ngrifos > resumos\nideias atomicas",
         "Etapa mais negligenciada — onde mora o valor.",
         PRIMARY_700),
        ("E", "EXPRESS", "Transformar em produto",
         "Drafts da dissertacao\nartigos\napresentacoes",
         "Output e o teste real do conhecimento.",
         SUCCESS),
    ]
    for i, (letter, name, sub, examples, insight, color) in enumerate(stages):
        x = Inches(1.2 + i * 2.95)
        rounded(s, x, Inches(3.4), Inches(2.8), Inches(2.65),
                NEUTRAL_50, line=NEUTRAL_200)
        bar(s, x, Inches(3.4), Inches(2.8), Pt(3), clr=color)
        rounded(s, x + Inches(0.2), Inches(3.55), Inches(0.45), Inches(0.45), color)
        txt(s, x + Inches(0.2), Inches(3.62), Inches(0.45), Inches(0.4),
            letter, FONT_HEADING, 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.75), Inches(3.55), Inches(2), Inches(0.3),
            name, FONT_HEADING, 13, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.75), Inches(3.85), Inches(2), Inches(0.25),
            sub, FONT_MONO, 9, color, bold=True)
        txt(s, x + Inches(0.2), Inches(4.3), Inches(2.6), Inches(0.95),
            examples, FONT_BODY, 10, NEUTRAL_700, spacing=14)
        # divider
        bar(s, x + Inches(0.2), Inches(5.35), Inches(2.4), Pt(0.75), clr=NEUTRAL_200)
        txt(s, x + Inches(0.2), Inches(5.45), Inches(2.6), Inches(0.6),
            insight, FONT_DISPLAY, 10, color, spacing=14)
        # arrow between cards
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + Inches(2.8), Inches(4.5),
                                      Inches(0.18), Inches(0.22))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT_500
            arr.line.fill.background()

    # Bottom: quando funciona / quando nao + recomendacao
    rounded(s, Inches(1.2), Inches(6.2), Inches(3.7), Inches(0.7),
            DEW_GREEN, line=SUCCESS)
    txt(s, Inches(1.4), Inches(6.27), Inches(3.5), Inches(0.25),
        "QUANDO FUNCIONA", FONT_HEADING, 9, SUCCESS, bold=True)
    txt(s, Inches(1.4), Inches(6.5), Inches(3.5), Inches(0.35),
        "Pesquisador com volume alto que precisa virar texto.",
        FONT_BODY, 10, NEUTRAL_700)

    rounded(s, Inches(5.1), Inches(6.2), Inches(3.7), Inches(0.7),
            DEW_GREEN, line=DANGER)
    txt(s, Inches(5.3), Inches(6.27), Inches(3.5), Inches(0.25),
        "QUANDO E LIMITADO", FONT_HEADING, 9, DANGER, bold=True)
    txt(s, Inches(5.3), Inches(6.5), Inches(3.5), Inches(0.35),
        "Sem habito de captura, o pipeline trava no inicio.",
        FONT_BODY, 10, NEUTRAL_700)

    rounded(s, Inches(9.0), Inches(6.2), Inches(3.7), Inches(0.7),
            PRIMARY_950)
    txt(s, Inches(9.2), Inches(6.27), Inches(3.5), Inches(0.25),
        "PARA O MESTRANDO", FONT_HEADING, 9, ACCENT_400, bold=True)
    txt(s, Inches(9.2), Inches(6.5), Inches(3.5), Inches(0.35),
        "Fluxo natural pesquisa > escrita. Cada agente = uma etapa.",
        FONT_BODY, 10, NEUTRAL_300)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 06. PLATAFORMA ----------------
def s06_plataforma(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "03  PLATAFORMA")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Google Drive como infraestrutura", FONT_HEADING, 32, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Trade-off analisado contra OneDrive. Decisao baseada em custo-beneficio + ecosistema academico.",
        FONT_BODY, 14, NEUTRAL_600)

    headers = ["Criterio", "Google Drive", "OneDrive", "Vencedor"]
    rows = [
        ("Espaco gratuito", "15 GB", "5 GB", "GDrive"),
        ("OCR em PDFs", "Nativo", "Limitado", "GDrive"),
        ("Integracao Scholar", "Library nativa", "Nao tem", "GDrive"),
        ("Google Colab (analise)", "Nativo", "Nao", "GDrive"),
        ("Word/Excel desktop", "Compativel", "Nativo", "OneDrive"),
        ("Plano EDU universidade", "Workspace EDU", "M365 EDU 1TB", "Depende"),
    ]
    col_w = [Inches(3.5), Inches(3.0), Inches(3.0), Inches(2.0)]
    col_x = [Inches(1.2), Inches(4.7), Inches(7.7), Inches(10.7)]
    y = Inches(2.55)
    filled_rect(s, Inches(1.2), y, Inches(11.5), Inches(0.5), PRIMARY_950)
    for i, h in enumerate(headers):
        txt(s, col_x[i], Inches(2.6), col_w[i], Inches(0.4),
            h.upper(), FONT_HEADING, 11, ACCENT_400, bold=True,
            align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)
        if i == 0:
            txt(s, col_x[i] + Inches(0.2), Inches(2.6), col_w[i], Inches(0.4),
                "CRITERIO", FONT_HEADING, 11, ACCENT_400, bold=True)
    for r, row in enumerate(rows):
        ry = Inches(3.05 + r * 0.45)
        if r % 2 == 0:
            filled_rect(s, Inches(1.2), ry, Inches(11.5), Inches(0.45), NEUTRAL_50)
        for i, val in enumerate(row):
            color = NEUTRAL_700
            bold = False
            font = FONT_BODY
            if i == 0:
                font = FONT_HEADING
                bold = True
                color = NEUTRAL_950
            if i == 3:
                font = FONT_MONO
                bold = True
                color = SUCCESS if val == "GDrive" else (PRIMARY_700 if val == "OneDrive" else NEUTRAL_500)
            txt(s, col_x[i] + (Inches(0.2) if i == 0 else Inches(0)),
                Inches(3.13 + r * 0.45), col_w[i], Inches(0.4),
                val, font, 12, color, bold=bold,
                align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    rounded(s, Inches(1.2), Inches(5.85), Inches(11.5), Inches(1.05), PRIMARY_950)
    bar(s, Inches(1.5), Inches(6.05), Inches(0.6))
    txt(s, Inches(1.5), Inches(6.18), Inches(11), Inches(0.4),
        "DECISAO  •  Google Drive como default", FONT_HEADING, 14, ACCENT_400, bold=True)
    txt(s, Inches(1.5), Inches(6.55), Inches(11), Inches(0.35),
        "Migre para OneDrive apenas se sua universidade oferece M365 EDU com 1TB e voce trabalha intensivo no Word desktop + EndNote.",
        FONT_BODY, 11, NEUTRAL_300)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 07. ESTRUTURA DE PASTAS ----------------
def s07_estrutura(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "04")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Arvore PARA numerada", FONT_HEADING, 32, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "Numeracao 00-99 garante ordem visual estavel no Drive.",
        FONT_BODY, 14, NEUTRAL_400)

    rounded(s, Inches(1.2), Inches(2.8), Inches(7.0), Inches(4.0),
            PRIMARY_900, line=PRIMARY_700, line_w=Pt(0.5))

    tree = [
        ("Mestrado-{Tema}/", 0, ACCENT_500, True),
        ("00_Inbox/", 1, DANGER, False),
        ("01_Projects/  (P do PARA: tem prazo)", 1, ACCENT_400, True),
        ("2026-Dissertacao/", 2, NEUTRAL_400, False),
        ("2026-Q3-Artigo-Periodico/", 2, NEUTRAL_400, False),
        ("02_Areas/  (A do PARA: continuas)", 1, ACCENT_400, True),
        ("Linha-Pesquisa/  Orientacao/  Disciplinas/", 2, NEUTRAL_400, False),
        ("03_Resources/  (R do PARA: tematicos)", 1, ACCENT_400, True),
        ("Bibliografia/  Notas-Atomicas/  Templates/", 2, NEUTRAL_400, False),
        ("04_Archive/  (concluido ou inativo)", 1, ACCENT_400, False),
        ("99_Sandbox/  (rascunhos descartaveis)", 1, NEUTRAL_500, False),
    ]
    y = 3.0
    for label, indent, color, bold in tree:
        x = Inches(1.45 + indent * 0.4)
        prefix = ("|-- " if indent > 0 else "")
        txt(s, x, Inches(y), Inches(7), Inches(0.32),
            prefix + label, FONT_MONO, 12, color, bold=bold)
        y += 0.32

    # Side note
    rounded(s, Inches(8.6), Inches(2.8), Inches(4.1), Inches(2.0),
            PRIMARY_900, line=ACCENT_500, line_w=Pt(0.75))
    bar(s, Inches(8.85), Inches(3.0), Inches(0.6))
    txt(s, Inches(8.85), Inches(3.15), Inches(3.7), Inches(0.4),
        "POR QUE 00 → 99", FONT_HEADING, 12, ACCENT_400, bold=True)
    txt(s, Inches(8.85), Inches(3.5), Inches(3.7), Inches(1.4),
        "Drive ordena lexicograficamente.\n\n00_Inbox sempre no topo —\nlembra de processar.\n\n99_Sandbox no fundo —\nisola experimentos.",
        FONT_BODY, 11, NEUTRAL_300, spacing=18)

    rounded(s, Inches(8.6), Inches(5.0), Inches(4.1), Inches(1.8),
            DEW_GREEN, line=ACCENT_500, line_w=Pt(0.75))
    bar(s, Inches(8.85), Inches(5.2), Inches(0.6), clr=PRIMARY_950)
    txt(s, Inches(8.85), Inches(5.35), Inches(3.7), Inches(0.4),
        "INSTALACAO", FONT_HEADING, 12, ACCENT_600, bold=True)
    txt(s, Inches(8.85), Inches(5.7), Inches(3.7), Inches(1.0),
        "Drive for Desktop\nsincroniza local.\nAgentes operam offline.\nMudancas sobem auto.",
        FONT_BODY, 11, NEUTRAL_700, spacing=18)
    footer(s, pg, TOTAL)


# ---------------- 08. NOMENCLATURA ----------------
def s08_nomenclatura(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "05  NOMENCLATURA")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Convencao YYYYMMDD-tipo-slug", FONT_HEADING, 32, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Data ISO ordenavel + tipo no nome = pesquisa em segundos.",
        FONT_BODY, 14, NEUTRAL_600)

    # Hero exemplo
    rounded(s, Inches(1.2), Inches(2.5), Inches(11.5), Inches(1.2),
            PRIMARY_950)
    txt(s, Inches(1.2), Inches(2.65), Inches(11.5), Inches(0.7),
        "20260507-lit-foucault-vigiar-punir.md",
        FONT_MONO, 26, ACCENT_500, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(1.2), Inches(3.2), Inches(11.5), Inches(0.4),
        "data ISO  •  tipo canonico  •  slug em kebab-case  •  extensao",
        FONT_BODY, 11, NEUTRAL_400, align=PP_ALIGN.CENTER)

    # Tipos canonicos grid 4x3
    types = [
        ("lit", "Fichamento de literatura"),
        ("nota", "Nota atomica permanente"),
        ("draft", "Rascunho de texto autoral"),
        ("dado", "Dataset bruto/tratado"),
        ("fig", "Figura, diagrama"),
        ("tab", "Tabela"),
        ("slides", "Apresentacao"),
        ("audio", "Audio (entrevista, voz)"),
        ("transcript", "Transcricao"),
        ("analise", "Output de analise"),
        ("peer", "Peer review"),
        ("meeting", "Ata de orientacao"),
    ]
    cols = 4
    for i, (prefix, desc) in enumerate(types):
        col = i % cols
        row = i // cols
        x = Inches(1.2 + col * 2.9)
        y = Inches(4.0 + row * 0.95)
        rounded(s, x, y, Inches(2.7), Inches(0.8), NEUTRAL_50,
                line=NEUTRAL_200)
        rounded(s, x, y, Inches(0.8), Inches(0.8), PRIMARY_950)
        txt(s, x, y + Inches(0.22), Inches(0.8), Inches(0.4),
            prefix, FONT_MONO, 14, ACCENT_400, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.95), y + Inches(0.25), Inches(1.65), Inches(0.4),
            desc, FONT_BODY, 10, NEUTRAL_700)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 09. WORKFLOW CODE ----------------
def s09_code(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "06")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Workflow CODE adaptado", FONT_HEADING, 32, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "Como o conteudo flui pelas pastas e pelos agentes.",
        FONT_BODY, 14, NEUTRAL_400)

    stages = [
        ("C", "CAPTURE", "00_Inbox/",
         "Tudo cai aqui sem julgamento.\nPDF, foto, audio, snippet.",
         "academic-researcher", DANGER),
        ("O", "ORGANIZE", "03_Resources/",
         "Move para PARA.\nPDFs, biblio.bib.",
         "citation-manager", PRIMARY_700),
        ("D", "DISTILL", "Notas-Atomicas/",
         "Uma ideia por nota.\nLinks bidirecionais.",
         "knowledge-architect", ACCENT_600),
        ("E", "EXPRESS", "01_Projects/drafts/",
         "Texto autoral.\nValidate > Audit > Humanize.",
         "academic-writer + crew", SUCCESS),
    ]
    for i, (letter, name, folder, desc, agent, color) in enumerate(stages):
        x = Inches(1.2 + i * 2.95)
        rounded(s, x, Inches(2.85), Inches(2.7), Inches(4.05),
                PRIMARY_900, line=color, line_w=Pt(1))
        # Big letter circle
        icon_circle(s, x + Inches(1.35), Inches(3.4), Inches(0.45),
                     color, letter, WHITE, mark_size=22, mark_font=FONT_HEADING)
        txt(s, x, Inches(4.0), Inches(2.7), Inches(0.4),
            name, FONT_HEADING, 14, ACCENT_500, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), Inches(4.45), Inches(2.3), Inches(0.3),
            folder, FONT_MONO, 10, NEUTRAL_300, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), Inches(4.85), Inches(2.3), Inches(0.9),
            desc, FONT_BODY, 11, NEUTRAL_300, spacing=16, align=PP_ALIGN.CENTER)
        # divider
        bar(s, x + Inches(0.6), Inches(5.95), Inches(1.5), Pt(1), clr=color)
        txt(s, x + Inches(0.2), Inches(6.05), Inches(2.3), Inches(0.25),
            "AGENTE", FONT_HEADING, 8, NEUTRAL_400, bold=True, align=PP_ALIGN.CENTER)
        txt(s, x + Inches(0.2), Inches(6.3), Inches(2.3), Inches(0.5),
            agent, FONT_MONO, 9, ACCENT_400, bold=True, align=PP_ALIGN.CENTER)
        # arrow between cards
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + Inches(2.7), Inches(4.7),
                                      Inches(0.25), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT_500
            arr.line.fill.background()
    footer(s, pg, TOTAL)


# ---------------- 10. SQUAD MAPA ----------------
def s10_squad(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "07  SQUAD")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "9 agentes especializados", FONT_HEADING, 32, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Cada agente atua sobre pastas especificas do Drive. Acionados na ordem do workflow.",
        FONT_BODY, 14, NEUTRAL_600)

    # Voce no centro
    rounded(s, Inches(5.4), Inches(2.5), Inches(2.5), Inches(0.7),
            PRIMARY_950)
    txt(s, Inches(5.4), Inches(2.6), Inches(2.5), Inches(0.3),
        "VOCE", FONT_HEADING, 14, ACCENT_500, bold=True, align=PP_ALIGN.CENTER)
    txt(s, Inches(5.4), Inches(2.9), Inches(2.5), Inches(0.25),
        "mestrando", FONT_MONO, 10, NEUTRAL_300, align=PP_ALIGN.CENTER)

    agents = [
        ("methodology-advisor",  "Desenho metodologico",  0.7,  3.5, PRIMARY_700),
        ("academic-researcher",  "Curadoria de fontes",   3.7,  3.5, PRIMARY_700),
        ("source-validator",     "Anti-alucinacao",       6.7,  3.5, DANGER),
        ("knowledge-architect",  "Zettelkasten",          9.7,  3.5, ACCENT_600),
        ("academic-writer",      "Redacao ABNT/APA",      0.7,  4.7, PRIMARY_950),
        ("citation-manager",     "Bibliografia .bib",     3.7,  4.7, NEUTRAL_700),
        ("ai-pattern-auditor",   "Mapa de calor de IA",   6.7,  4.7, DANGER),
        ("voice-humanizer",      "Voz autoral",           9.7,  4.7, ACCENT_600),
        ("peer-reviewer",        "Banca simulada",       4.7,   5.9, PRIMARY_950),
    ]
    for name, role, x_in, y_in, color in agents:
        x = Inches(x_in)
        y = Inches(y_in)
        rounded(s, x, y, Inches(2.95), Inches(0.95), DEW_GREEN,
                line=NEUTRAL_200, line_w=Pt(0.75))
        bar(s, x, y, Inches(2.95), Pt(2.5), clr=color)
        txt(s, x + Inches(0.15), y + Inches(0.15), Inches(2.7), Inches(0.35),
            name, FONT_MONO, 11, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.15), y + Inches(0.5), Inches(2.7), Inches(0.35),
            role, FONT_BODY, 11, NEUTRAL_700)
    # Highlight novo
    txt(s, Inches(6.7), Inches(4.7) - Inches(0.3), Inches(2.95), Inches(0.25),
        "NOVO", FONT_MONO, 9, ACCENT_600, bold=True, align=PP_ALIGN.RIGHT)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 11. AGENTES PESQUISA + VALIDACAO ----------------
def s11_agentes_a(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "07.A")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Pesquisa & curadoria", FONT_HEADING, 30, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "methodology-advisor  →  academic-researcher  →  source-validator",
        FONT_MONO, 13, ACCENT_400)

    agents = [
        {
            "name": "methodology-advisor", "role": "Desenho metodologico",
            "le": ["02_Areas/Orientacao/", "memory/project_thesis.md"],
            "esc": ["01_Projects/{atual}/metodologia.md", "memory/project_methodology.md"],
            "trigger": "Inicio de cada novo projeto.",
        },
        {
            "name": "academic-researcher", "role": "Curadoria de fontes",
            "le": ["memory/project_thesis.md", "memory/project_keywords.md"],
            "esc": ["03_Resources/Bibliografia/PDFs/", "research/{data}-{slug}/sources.md"],
            "trigger": "Revisao de literatura e atualizacao do estado da arte.",
        },
        {
            "name": "source-validator", "role": "Anti-alucinacao",
            "le": ["drafts/*.md", "biblio.bib"],
            "esc": ["peer-review/{data}-validacao-{slug}.md"],
            "trigger": "APOS qualquer output de IA, ANTES de submissao.",
        },
    ]
    for i, a in enumerate(agents):
        x = Inches(0.6 + i * 4.15)
        rounded(s, x, Inches(2.85), Inches(3.95), Inches(4.0),
                PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
        bar(s, x + Inches(0.3), Inches(3.05), Inches(0.5))
        txt(s, x + Inches(0.3), Inches(3.2), Inches(3.5), Inches(0.4),
            a["name"], FONT_MONO, 13, ACCENT_500, bold=True)
        txt(s, x + Inches(0.3), Inches(3.55), Inches(3.5), Inches(0.3),
            a["role"], FONT_BODY, 11, NEUTRAL_400)
        txt(s, x + Inches(0.3), Inches(4.0), Inches(3.5), Inches(0.3),
            "LE", FONT_HEADING, 9, ACCENT_400, bold=True)
        for j, item in enumerate(a["le"]):
            txt(s, x + Inches(0.3), Inches(4.25 + j*0.28), Inches(3.5), Inches(0.3),
                "• " + item, FONT_MONO, 9, NEUTRAL_300)
        ye = 4.25 + len(a["le"]) * 0.28 + 0.1
        txt(s, x + Inches(0.3), Inches(ye), Inches(3.5), Inches(0.3),
            "ESCREVE", FONT_HEADING, 9, ACCENT_400, bold=True)
        for j, item in enumerate(a["esc"]):
            txt(s, x + Inches(0.3), Inches(ye + 0.25 + j*0.28), Inches(3.5), Inches(0.3),
                "• " + item, FONT_MONO, 9, NEUTRAL_300)
        txt(s, x + Inches(0.3), Inches(6.2), Inches(3.5), Inches(0.3),
            "TRIGGER", FONT_HEADING, 9, ACCENT_400, bold=True)
        txt(s, x + Inches(0.3), Inches(6.45), Inches(3.5), Inches(0.4),
            a["trigger"], FONT_BODY, 10, NEUTRAL_300, spacing=14)
    footer(s, pg, TOTAL)


# ---------------- 12. AGENTES CONHECIMENTO + REDACAO ----------------
def s12_agentes_b(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "07.B  CONHECIMENTO & REDACAO")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Conhecimento & redacao", FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "knowledge-architect  →  academic-writer  →  citation-manager",
        FONT_MONO, 13, ACCENT_600)

    agents = [
        {
            "name": "knowledge-architect", "role": "Zettelkasten",
            "le": ["PDFs em Bibliografia/", "Audios dictados"],
            "esc": ["Notas-Atomicas/permanent/{ts}-{slug}.md", "Notas-Atomicas/index/moc-{tema}.md"],
            "trigger": "A cada leitura/insight + rito quinzenal de promocao.",
        },
        {
            "name": "academic-writer", "role": "Redacao ABNT/APA",
            "le": ["Notas-Atomicas/permanent/", "biblio.bib", "memory/user_writing_voice.md"],
            "esc": ["01_Projects/{atual}/drafts/{capitulo}.md"],
            "trigger": "Sprint de redacao com 8-12 notas conectadas.",
        },
        {
            "name": "citation-manager", "role": "Bibliografia .bib",
            "le": ["Listas de referencias", "DOIs"],
            "esc": ["biblio.bib (UNICO no Drive)", "referencias-{norma}.md"],
            "trigger": "Toda nova fonte + antes de cada submissao.",
        },
    ]
    for i, a in enumerate(agents):
        x = Inches(0.6 + i * 4.15)
        rounded(s, x, Inches(2.85), Inches(3.95), Inches(4.0),
                DEW_GREEN, line=ACCENT_500, line_w=Pt(0.5))
        bar(s, x + Inches(0.3), Inches(3.05), Inches(0.5))
        txt(s, x + Inches(0.3), Inches(3.2), Inches(3.5), Inches(0.4),
            a["name"], FONT_MONO, 13, ACCENT_600, bold=True)
        txt(s, x + Inches(0.3), Inches(3.55), Inches(3.5), Inches(0.3),
            a["role"], FONT_BODY, 11, NEUTRAL_700)
        txt(s, x + Inches(0.3), Inches(4.0), Inches(3.5), Inches(0.3),
            "LE", FONT_HEADING, 9, NEUTRAL_950, bold=True)
        for j, item in enumerate(a["le"]):
            txt(s, x + Inches(0.3), Inches(4.25 + j*0.28), Inches(3.5), Inches(0.3),
                "• " + item, FONT_MONO, 9, NEUTRAL_700)
        ye = 4.25 + len(a["le"]) * 0.28 + 0.1
        txt(s, x + Inches(0.3), Inches(ye), Inches(3.5), Inches(0.3),
            "ESCREVE", FONT_HEADING, 9, NEUTRAL_950, bold=True)
        for j, item in enumerate(a["esc"]):
            txt(s, x + Inches(0.3), Inches(ye + 0.25 + j*0.28), Inches(3.5), Inches(0.3),
                "• " + item, FONT_MONO, 9, NEUTRAL_700)
        txt(s, x + Inches(0.3), Inches(6.2), Inches(3.5), Inches(0.3),
            "TRIGGER", FONT_HEADING, 9, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.3), Inches(6.45), Inches(3.5), Inches(0.4),
            a["trigger"], FONT_BODY, 10, NEUTRAL_700, spacing=14)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 13. AGENTES POLIMENTO ----------------
def s13_agentes_c(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "07.C")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Polimento & banca", FONT_HEADING, 30, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "ai-pattern-auditor  →  voice-humanizer  →  peer-reviewer",
        FONT_MONO, 13, ACCENT_400)

    agents = [
        {
            "name": "ai-pattern-auditor", "role": "Mapa de calor de IA",
            "tag": "NOVO",
            "le": ["memory/user_writing_voice.md", "drafts a auditar"],
            "esc": ["peer-review/{data}-ai-audit-{slug}.md", "memory/ai_audit_log.md"],
            "trigger": "ANTES do voice-humanizer. Aponta trechos para reescrita focada.",
        },
        {
            "name": "voice-humanizer", "role": "Voz autoral",
            "tag": "QUEST",
            "le": ["memory/voice_questionnaire.md", "user_writing_voice.md"],
            "esc": ["drafts reescritos com aluno", "memory/voice_calibration_log.md"],
            "trigger": "Aplica QUESTIONARIO de calibracao na 1a sessao.",
        },
        {
            "name": "peer-reviewer", "role": "Banca simulada",
            "tag": "",
            "le": ["drafts completos", "metodologia.md", "MOCs"],
            "esc": ["peer-review/{data}-banca-{slug}.md", "memory/banca_questions.md"],
            "trigger": "ANTES de qualificacao, defesa ou submissao.",
        },
    ]
    for i, a in enumerate(agents):
        x = Inches(0.6 + i * 4.15)
        rounded(s, x, Inches(2.85), Inches(3.95), Inches(4.0),
                PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
        if a["tag"]:
            tag_box = rounded(s, x + Inches(2.95), Inches(2.7),
                               Inches(0.85), Inches(0.3), ACCENT_500)
            txt(s, x + Inches(2.95), Inches(2.74), Inches(0.85), Inches(0.25),
                a["tag"], FONT_MONO, 9, PRIMARY_950, bold=True,
                align=PP_ALIGN.CENTER)
        bar(s, x + Inches(0.3), Inches(3.05), Inches(0.5))
        txt(s, x + Inches(0.3), Inches(3.2), Inches(3.5), Inches(0.4),
            a["name"], FONT_MONO, 13, ACCENT_500, bold=True)
        txt(s, x + Inches(0.3), Inches(3.55), Inches(3.5), Inches(0.3),
            a["role"], FONT_BODY, 11, NEUTRAL_400)
        txt(s, x + Inches(0.3), Inches(4.0), Inches(3.5), Inches(0.3),
            "LE", FONT_HEADING, 9, ACCENT_400, bold=True)
        for j, item in enumerate(a["le"]):
            txt(s, x + Inches(0.3), Inches(4.25 + j*0.28), Inches(3.5), Inches(0.3),
                "• " + item, FONT_MONO, 9, NEUTRAL_300)
        ye = 4.25 + len(a["le"]) * 0.28 + 0.1
        txt(s, x + Inches(0.3), Inches(ye), Inches(3.5), Inches(0.3),
            "ESCREVE", FONT_HEADING, 9, ACCENT_400, bold=True)
        for j, item in enumerate(a["esc"]):
            txt(s, x + Inches(0.3), Inches(ye + 0.25 + j*0.28), Inches(3.5), Inches(0.3),
                "• " + item, FONT_MONO, 9, NEUTRAL_300)
        txt(s, x + Inches(0.3), Inches(6.2), Inches(3.5), Inches(0.3),
            "TRIGGER", FONT_HEADING, 9, ACCENT_400, bold=True)
        txt(s, x + Inches(0.3), Inches(6.45), Inches(3.5), Inches(0.5),
            a["trigger"], FONT_BODY, 10, NEUTRAL_300, spacing=14)
    footer(s, pg, TOTAL)


# ---------------- 14. VOICE QUESTIONARIO ----------------
def s14_voice_questionario(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "08  VOICE-HUMANIZER")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Questionario de calibracao", FONT_HEADING, 32, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "21 perguntas em 5 blocos. Aplicado na primeira sessao com cada aluno.",
        FONT_BODY, 14, NEUTRAL_600)

    blocks = [
        ("Bloco 1", "Identidade academica", "5 min",
         "Area, nivel de formalidade, eu/nos/impessoal,\nautores admirados/rejeitados",
         PRIMARY_950),
        ("Bloco 2", "Repertorio linguistico", "10 min",
         "5 palavras que voce usa, 5 que rejeita,\nfrase curta/longa, verbo padrao de posicionamento",
         PRIMARY_700),
        ("Bloco 3", "Habitos de escrita", "10 min",
         "Direto/oral/manual, como reage ao bloqueio,\ntempo em revisao, sensacao ao reler",
         ACCENT_600),
        ("Bloco 4", "Posicionamento e ethos", "5 min",
         "Conciliador/provocador/didatico/sintetico,\ncomo aborda autor citado, tema com urgencia",
         DANGER),
        ("Bloco 5", "Materiais para analise", "5 min",
         "3 textos seus bons + 1 odiado + 1 informal\npara extracao de padroes estilisticos",
         NEUTRAL_700),
    ]
    for i, (b, t, time, desc, color) in enumerate(blocks):
        y = Inches(2.6 + i * 0.85)
        rounded(s, Inches(1.2), y, Inches(11.5), Inches(0.78),
                NEUTRAL_50, line=NEUTRAL_200)
        rounded(s, Inches(1.2), y, Inches(1.9), Inches(0.78), color)
        txt(s, Inches(1.2), y + Inches(0.12), Inches(1.9), Inches(0.3),
            b, FONT_HEADING, 13, ACCENT_400 if color != ACCENT_600 else PRIMARY_950,
            bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(1.2), y + Inches(0.42), Inches(1.9), Inches(0.25),
            time, FONT_MONO, 10,
            NEUTRAL_300 if color != ACCENT_600 else PRIMARY_950,
            align=PP_ALIGN.CENTER)
        txt(s, Inches(3.3), y + Inches(0.07), Inches(9.2), Inches(0.35),
            t, FONT_HEADING, 14, NEUTRAL_950, bold=True)
        txt(s, Inches(3.3), y + Inches(0.4), Inches(9.2), Inches(0.4),
            desc, FONT_BODY, 11, NEUTRAL_700, spacing=14)
    rounded(s, Inches(1.2), Inches(6.95), Inches(11.5), Inches(0.42),
            PRIMARY_950)
    txt(s, Inches(1.4), Inches(7.02), Inches(11), Inches(0.3),
        "SAIDA  •  memory/voice_questionnaire.md (bruto) + memory/user_writing_voice.md (perfil destilado)",
        FONT_MONO, 11, ACCENT_400, bold=True)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 15. AI AUDITOR HEATMAP ----------------
def s15_auditor(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "08.B")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "AI Pattern Auditor", FONT_HEADING, 32, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "Diagnostica trechos com marca de IA. NAO reescreve — entrega mapa de calor para humanizacao focada.",
        FONT_BODY, 14, NEUTRAL_400)

    # Esquerda: taxonomia de marcadores
    rounded(s, Inches(1.2), Inches(2.85), Inches(5.7), Inches(4.0),
            PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
    txt(s, Inches(1.4), Inches(3.0), Inches(5.3), Inches(0.4),
        "TAXONOMIA DE MARCADORES", FONT_HEADING, 11, ACCENT_500, bold=True)
    cats = [
        ("A", "Estatisticos", "burstiness, comprimento, repeticao", PRIMARY_700),
        ("B", "Lexicais", "clichês, hedging, tripartites", ACCENT_600),
        ("C", "Retoricos", "sem posicionamento, conclusao redundante", DANGER),
        ("D", "Voz ausente", "verbo errado, sem palavra-forte do aluno", SUCCESS),
        ("E", "Integridade", "citacao perfeita demais, dado sem fonte", WARN),
    ]
    for i, (code, name, desc, color) in enumerate(cats):
        y = Inches(3.45 + i * 0.6)
        # code badge
        icon_circle(s, Inches(1.7), y + Inches(0.2), Inches(0.18),
                     color, code, WHITE, mark_size=11, mark_font=FONT_HEADING)
        txt(s, Inches(2.05), y + Inches(0.04), Inches(2.0), Inches(0.3),
            name, FONT_HEADING, 13, WHITE, bold=True)
        txt(s, Inches(2.05), y + Inches(0.3), Inches(4.5), Inches(0.3),
            desc, FONT_BODY, 10, NEUTRAL_400)

    # Direita: heatmap exemplo
    rounded(s, Inches(7.1), Inches(2.85), Inches(5.5), Inches(4.0),
            PRIMARY_900, line=ACCENT_500, line_w=Pt(0.5))
    txt(s, Inches(7.3), Inches(3.0), Inches(5.1), Inches(0.4),
        "EXEMPLO DE HEATMAP", FONT_HEADING, 11, ACCENT_500, bold=True)

    rows = [
        ("P1", 8.7, "CRITICO", DANGER, "B1, B3, C1, D2, D6 — reescrever"),
        ("P2", 6.4, "ALTO",    WARN, "B7, C7, D3 — cirurgia"),
        ("P3", 3.1, "BAIXO",   SUCCESS, "voz consolidada — manter"),
        ("P4", 8.2, "CRITICO", DANGER, "B1, C3, D4 — reescrever"),
        ("P5", 5.5, "MEDIO",   WARN, "B4 — ajuste pontual"),
    ]
    for i, (p, score, label, color, ann) in enumerate(rows):
        y = Inches(3.5 + i * 0.6)
        # paragraph badge
        rounded(s, Inches(7.3), y, Inches(0.5), Inches(0.45),
                PRIMARY_950, line=color, line_w=Pt(1))
        txt(s, Inches(7.3), y + Inches(0.07), Inches(0.5), Inches(0.3),
            p, FONT_MONO, 11, ACCENT_400, bold=True, align=PP_ALIGN.CENTER)
        # score bar
        bar_w = Inches(score * 0.32)
        filled_rect(s, Inches(7.95), y + Inches(0.18), Inches(3.2), Pt(2),
                     PRIMARY_950)
        filled_rect(s, Inches(7.95), y + Inches(0.18), bar_w, Pt(2), color)
        txt(s, Inches(7.95), y, Inches(3.2), Inches(0.2),
            f"{score:.1f} • {label}", FONT_MONO, 9, color, bold=True)
        txt(s, Inches(7.95), y + Inches(0.25), Inches(4.6), Inches(0.25),
            ann, FONT_BODY, 9, NEUTRAL_300)
    footer(s, pg, TOTAL)


# ---------------- 16. PESQUISA RAPIDA ----------------
def s16_busca(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, WHITE)
    left_bar(s)
    section_label(s, "09  PESQUISA RAPIDA")
    add_logo(s, "horiz_dark", Inches(11.5), Inches(0.5), Inches(0.4))
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Achar qualquer coisa em 5 segundos", FONT_HEADING, 30, NEUTRAL_950, bold=True)
    txt(s, Inches(1.2), Inches(1.85), Inches(11), Inches(0.4),
        "Cinco estrategias de busca cobrindo todos os cenarios da jornada.",
        FONT_BODY, 14, NEUTRAL_600)

    strategies = [
        ("01", "Drive Web",
         "Interface oficial",
         ["type:pdf foucault", "before:2026-01-01", "OCR busca dentro do PDF"],
         PRIMARY_950),
        ("02", "Windows Search",
         "Drive for Desktop",
         ["Ctrl+F no Explorer", "Filtro ISO ordenavel", "Pre-visualizacao"],
         PRIMARY_700),
        ("03", "Busca semantica",
         "knowledge-architect",
         ["Liste notas Foucault↔Deleuze", "Por relacao semantica", "Cobre o que tag nao captura"],
         ACCENT_600),
        ("04", "Tag (grep)",
         "Frontmatter YAML",
         ["grep -lr 'tags:.*foucault'", "Filtro por status", "Listagem por projeto"],
         DANGER),
        ("05", "MOCs",
         "Curado manualmente",
         ["index/moc-{tema}.md", "Portao de entrada", "Comece pelo MOC"],
         NEUTRAL_700),
    ]
    for i, (num, name, sub, items, color) in enumerate(strategies):
        col = i % 3
        row = i // 3
        x = Inches(0.7 + col * 4.2)
        y = Inches(2.55 + row * 2.3)
        rounded(s, x, y, Inches(4.0), Inches(2.15), DEW_GREEN,
                line=NEUTRAL_200)
        bar(s, x, y, Inches(4.0), Pt(3), clr=color)
        txt(s, x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.3),
            num, FONT_MONO, 12, color, bold=True)
        txt(s, x + Inches(0.95), y + Inches(0.15), Inches(2.7), Inches(0.3),
            name, FONT_HEADING, 14, NEUTRAL_950, bold=True)
        txt(s, x + Inches(0.2), y + Inches(0.5), Inches(3.7), Inches(0.25),
            sub, FONT_BODY, 10, NEUTRAL_600)
        for j, item in enumerate(items):
            yy = y + Inches(0.85 + j * 0.35)
            d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.25), yy + Inches(0.1),
                                    Inches(0.07), Inches(0.07))
            d.fill.solid()
            d.fill.fore_color.rgb = color
            d.line.fill.background()
            txt(s, x + Inches(0.45), yy, Inches(3.5), Inches(0.3),
                item, FONT_MONO, 10, NEUTRAL_700)
    footer(s, pg, TOTAL, dark=False)


# ---------------- 16b. INTERFACE OBSIDIAN ----------------
def s16b_obsidian(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "10")
    txt(s, Inches(1.2), Inches(1.05), Inches(11), Inches(0.7),
        "Obsidian — interface unificada do squad",
        FONT_HEADING, 30, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(1.85), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.0), Inches(11), Inches(0.4),
        "Tier 1 da estrategia de UI  •  gratuito, local, offline  •  30 min de setup",
        FONT_BODY, 13, NEUTRAL_400)

    # Hero
    rounded(s, Inches(1.2), Inches(2.6), Inches(11.5), Inches(0.85),
            PRIMARY_900, line=ACCENT_500)
    txt(s, Inches(1.5), Inches(2.75), Inches(11), Inches(0.3),
        "PRINCIPIO", FONT_HEADING, 10, ACCENT_500, bold=True)
    txt(s, Inches(1.5), Inches(3.03), Inches(11), Inches(0.4),
        "Mesmos arquivos .md que os agentes editam — sem duplicacao, sem sincronizacao.",
        FONT_DISPLAY, 14, WHITE, bold=True)

    # 4 superpoderes em cards 2x2
    powers = [
        ("01", "Grafo Zettelkasten",
         "Visualizacao da rede\nde notas linkadas",
         "Ctrl+G abre o grafo.\nFiltros por tag/pasta\nrevelam sub-redes.",
         ACCENT_600),
        ("02", "Dataview",
         "Queries SQL-like sobre\no frontmatter dos .md",
         'TABLE status, tags\nFROM "Notas-Atomicas"\nWHERE contains(tags,\n  "foucault")',
         PRIMARY_700),
        ("03", "Smart Connections",
         "Chat com seu vault\n+ busca semantica",
         '"Quais notas conectam\nFoucault e Deleuze?"\n\nResponde citando\nas notas usadas.',
         DANGER),
        ("04", "Templater",
         "Templates de fichamento,\nnota atomica, daily",
         "Ctrl+Shift+N\ncria nota atomica\ncom frontmatter\npre-preenchido.",
         SUCCESS),
    ]
    for i, (num, name, sub, demo, color) in enumerate(powers):
        col = i % 2
        row = i // 2
        x = Inches(1.2 + col * 5.85)
        y = Inches(3.7 + row * 1.65)
        rounded(s, x, y, Inches(5.6), Inches(1.5),
                PRIMARY_900, line=color, line_w=Pt(0.75))
        bar(s, x + Inches(0.3), y + Inches(0.2), Inches(0.5), clr=color)
        txt(s, x + Inches(0.3), y + Inches(0.32), Inches(0.5), Inches(0.3),
            num, FONT_MONO, 10, color, bold=True)
        txt(s, x + Inches(0.85), y + Inches(0.18), Inches(2.4), Inches(0.4),
            name, FONT_HEADING, 14, WHITE, bold=True)
        txt(s, x + Inches(0.85), y + Inches(0.55), Inches(2.4), Inches(0.7),
            sub, FONT_BODY, 10, NEUTRAL_400, spacing=14)
        # demo box (right side of card)
        rounded(s, x + Inches(3.4), y + Inches(0.18), Inches(2.05), Inches(1.2),
                PRIMARY_950, line=NEUTRAL_700)
        txt(s, x + Inches(3.55), y + Inches(0.28), Inches(1.85), Inches(1.0),
            demo, FONT_MONO, 9, ACCENT_400, spacing=14)

    # Bottom: integracao com squad
    rounded(s, Inches(1.2), Inches(7.0), Inches(11.5), Inches(0.32),
            ACCENT_500)
    txt(s, Inches(1.4), Inches(7.05), Inches(11.3), Inches(0.25),
        "INTEGRACAO COM O SQUAD  •  Agentes editam .md > Obsidian recarrega auto > Smart Connections re-embedda em background",
        FONT_MONO, 9, PRIMARY_950, bold=True)

    footer(s, pg, TOTAL)


# ---------------- 17. CRONOGRAMA ----------------
def s17_cronograma(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    chapter_num(s, "09")
    txt(s, Inches(1.2), Inches(1.1), Inches(11), Inches(0.7),
        "Cronograma de adocao", FONT_HEADING, 32, WHITE, bold=True)
    bar(s, Inches(1.2), Inches(2.0), Inches(1.2))
    txt(s, Inches(1.2), Inches(2.2), Inches(11), Inches(0.4),
        "De zero ao primeiro fluxo CODE completo em 4 semanas.",
        FONT_BODY, 14, NEUTRAL_400)

    weeks = [
        ("S1", "Estrutura", PRIMARY_700,
         ["Criar Mestrado-{Tema}/", "Replicar arvore PARA",
          "Drive for Desktop", "Mover material para Inbox"]),
        ("S2", "Calibracao", ACCENT_600,
         ["memory/project_thesis.md", "memory/citation_style.md",
          "Questionario voice-humanizer", "3 textos antigos para extracao"]),
        ("S3", "Primeiro CODE", DANGER,
         ["5 fontes via researcher", "Validar com source-validator",
          "Fichar com architect", "5 notas atomicas linkadas"]),
        ("S4", "Primeiro Express", SUCCESS,
         ["8-12 notas conectadas", "writer redige 1500 palavras",
          "Auditor + humanizer + peer", "Avaliar e ajustar"]),
    ]
    for i, (week, title, color, items) in enumerate(weeks):
        x = Inches(0.6 + i * 3.1)
        rounded(s, x, Inches(2.85), Inches(2.85), Inches(4.05),
                PRIMARY_900, line=color, line_w=Pt(1))
        # week number circle
        icon_circle(s, x + Inches(1.425), Inches(3.4),
                     Inches(0.4), color, week, WHITE,
                     mark_size=14, mark_font=FONT_HEADING)
        txt(s, x, Inches(3.95), Inches(2.85), Inches(0.4),
            title, FONT_HEADING, 18, ACCENT_500, bold=True, align=PP_ALIGN.CENTER)
        bar(s, x + Inches(1.0), Inches(4.4), Inches(0.85), Pt(1.5), clr=color)
        for j, item in enumerate(items):
            yy = Inches(4.65 + j * 0.45)
            d = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.25), yy + Inches(0.1),
                                    Inches(0.08), Inches(0.08))
            d.fill.solid()
            d.fill.fore_color.rgb = color
            d.line.fill.background()
            txt(s, x + Inches(0.45), yy, Inches(2.4), Inches(0.4),
                item, FONT_BODY, 10, NEUTRAL_300, spacing=14)
        # arrow
        if i < 3:
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                      x + Inches(2.85), Inches(4.7),
                                      Inches(0.25), Inches(0.3))
            arr.fill.solid()
            arr.fill.fore_color.rgb = ACCENT_500
            arr.line.fill.background()
    footer(s, pg, TOTAL)


# ---------------- 18. CLOSING / PROXIMOS PASSOS ----------------
def s18_closing(pg):
    s = prs.slides.add_slide(BLANK)
    set_bg(s, PRIMARY_950)
    filled_rect(s, Inches(0), Inches(2.6), SLIDE_W, Inches(0.04), ACCENT_500)
    add_logo(s, "horiz_light", Inches(10.5), Inches(0.5), Inches(0.45))

    txt(s, Inches(1.2), Inches(1.0), Inches(11), Inches(0.4),
        "PROXIMOS PASSOS", FONT_HEADING, 12, ACCENT_400, bold=True)
    txt(s, Inches(1.2), Inches(1.4), Inches(11), Inches(0.9),
        "Inteligencia que se constroi.", FONT_DISPLAY, 42, WHITE, bold=True)
    txt(s, Inches(1.2), Inches(2.15), Inches(11), Inches(0.35),
        "Da estrutura ao primeiro draft em 30 dias.",
        FONT_DISPLAY, 16, ACCENT_400)

    steps = [
        ("01", "Aprovar este plano",
         "Revisar arvore, nomenclatura e agentes — ajustar ao seu tema."),
        ("02", "Provisionar o Drive",
         "Criar Mestrado-{Tema}/ com toda a arvore PARA. Drive for Desktop."),
        ("03", "Aplicar questionario",
         "voice-humanizer aplica os 21 itens. Voce envia 3 textos antigos."),
        ("04", "Rodar Sprint 1",
         "5 fontes → validar → fichar → 5 notas atomicas linkadas."),
        ("05", "Avaliar e iterar",
         "Apos primeira mini-secao, ajustar prompts dos agentes que falharem."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.95 + i * 0.7)
        rounded(s, Inches(1.2), y, Inches(0.6), Inches(0.55), ACCENT_500)
        txt(s, Inches(1.2), y + Inches(0.13), Inches(0.6), Inches(0.4),
            num, FONT_MONO, 14, PRIMARY_950, bold=True, align=PP_ALIGN.CENTER)
        txt(s, Inches(2.0), y + Inches(0.05), Inches(11), Inches(0.35),
            title, FONT_HEADING, 16, WHITE, bold=True)
        txt(s, Inches(2.0), y + Inches(0.35), Inches(11), Inches(0.3),
            desc, FONT_BODY, 11, NEUTRAL_300)

    txt(s, Inches(1.2), Inches(6.85), Inches(11), Inches(0.3),
        "docs/curadoria-conteudo-plano.md  •  agents/  •  memory/  •  scripts/",
        FONT_MONO, 9, NEUTRAL_500)
    txt(s, Inches(1.2), Inches(7.15), Inches(11), Inches(0.3),
        "Verticalis AI Journey  •  Education  •  Maio 2026  •  v1.0",
        FONT_MONO, 9, ACCENT_400, align=PP_ALIGN.RIGHT)


# ============================================================================
# build all
# ============================================================================
s01_cover()
s02_toc(2)
s03_manifesto(3)
s04_diagnostico(4)
s05_framework(5)
s05a_para(6)
s05b_zettel(7)
s05c_code_detailed(8)
s06_plataforma(9)
s07_estrutura(10)
s08_nomenclatura(11)
s09_code(12)
s10_squad(13)
s11_agentes_a(14)
s12_agentes_b(15)
s13_agentes_c(16)
s14_voice_questionario(17)
s15_auditor(18)
s16_busca(19)
s16b_obsidian(20)
s17_cronograma(21)
s18_closing(22)

out = ROOT / "docs" / "curadoria-verticalis-deck.pptx"
prs.save(str(out))
print(f"PPTX gerado: {out}")
print(f"Total de slides: {len(prs.slides)}")
