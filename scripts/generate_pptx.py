"""
Gera apresentação PPTX do framework de curadoria acadêmica.
Uso: python scripts/generate_pptx.py
Saída: docs/curadoria-conteudo-framework.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------- paleta -----------------------
NAVY = RGBColor(0x0E, 0x2A, 0x47)
TEAL = RGBColor(0x1F, 0x7A, 0x8C)
GOLD = RGBColor(0xE8, 0xB4, 0x1B)
LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
GRAY = RGBColor(0x6B, 0x73, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x27, 0xAE, 0x60)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background() if line is None else None
    if line is not None:
        shape.line.color.rgb = line
    return shape


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, size=14, color=DARK, line_spacing=1.2):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def header_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), NAVY)
    add_text(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.6),
             title, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.55), Inches(12), Inches(0.35),
                 subtitle, size=12, color=LIGHT)
    add_rect(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), GOLD)


def footer_bar(slide, page_num, total):
    add_rect(slide, Inches(0), Inches(7.2), prs.slide_width, Inches(0.3), LIGHT)
    add_text(slide, Inches(0.5), Inches(7.22), Inches(8), Inches(0.3),
             "Sistema de Curadoria Academica - Squad de Agentes", size=9, color=GRAY)
    add_text(slide, Inches(11.5), Inches(7.22), Inches(1.3), Inches(0.3),
             f"{page_num} / {total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)


# ----------------------- 1. Capa -----------------------
TOTAL = 16


def slide_capa():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(0), Inches(3.2), prs.slide_width, Inches(1.8), TEAL)
    add_rect(s, Inches(0), Inches(5.0), prs.slide_width, Inches(0.08), GOLD)
    add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(1.2),
             "Sistema de Curadoria de Conteudo Academico",
             size=44, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(3.45), Inches(12), Inches(0.6),
             "PARA + Zettelkasten + CODE | Google Drive + Squad de 8 Agentes",
             size=20, color=WHITE)
    add_text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.5),
             "Pesquisa rapida, contexto persistente, autoria autentica",
             size=16, color=LIGHT)
    add_text(s, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
             "Mestrado | 2026-05-07 | v1.0",
             size=12, color=GOLD)


# ----------------------- 2. Problema -----------------------
def slide_problema():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "O problema que estamos resolvendo",
               "Tres frictions tipicas na pesquisa academica assistida por IA")
    add_rect(s, Inches(0.5), Inches(1.2), Inches(4), Inches(5.8), LIGHT)
    add_rect(s, Inches(4.7), Inches(1.2), Inches(4), Inches(5.8), LIGHT)
    add_rect(s, Inches(8.9), Inches(1.2), Inches(4), Inches(5.8), LIGHT)
    add_rect(s, Inches(0.5), Inches(1.2), Inches(4), Inches(0.6), RED)
    add_rect(s, Inches(4.7), Inches(1.2), Inches(4), Inches(0.6), RED)
    add_rect(s, Inches(8.9), Inches(1.2), Inches(4), Inches(0.6), RED)
    titles = ["Curadoria pulverizada", "Conhecimento se perde", "Texto soa de IA"]
    for i, t in enumerate(titles):
        add_text(s, Inches(0.5 + i*4.2), Inches(1.32), Inches(4), Inches(0.4),
                 t, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sintomas = [
        "Resultados de IA misturam fontes confiaveis, blogs e alucinacoes. Sem gatekeeper que valide cada afirmacao.",
        "Insights ficam em conversas dispersas, sem estrutura cumulativa. Cada sessao recomeca do zero.",
        "Detectores acusam padrao estatistico (baixa burstiness, vocabulario previsivel). Texto sem voz autoral."
    ]
    for i, txt in enumerate(sintomas):
        add_text(s, Inches(0.7 + i*4.2), Inches(2.0), Inches(3.6), Inches(2.5),
                 "Sintoma", size=11, bold=True, color=GRAY)
        add_text(s, Inches(0.7 + i*4.2), Inches(2.3), Inches(3.6), Inches(2.5),
                 txt, size=12, color=DARK)
    causas = [
        "Falta de fluxo Capture > Validate > Distill antes de armazenar.",
        "Falta de camada Zettelkasten / segundo cerebro persistente.",
        "Texto gerado em uma unica passada sem reescrita autoral ancorada na voz do aluno."
    ]
    for i, txt in enumerate(causas):
        add_text(s, Inches(0.7 + i*4.2), Inches(4.6), Inches(3.6), Inches(0.4),
                 "Causa raiz", size=11, bold=True, color=GRAY)
        add_text(s, Inches(0.7 + i*4.2), Inches(4.9), Inches(3.6), Inches(2),
                 txt, size=12, color=DARK)
    footer_bar(s, 2, TOTAL)


# ----------------------- 3. Visao geral -----------------------
def slide_visao():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Visao geral - 3 frameworks combinados",
               "Cada um cobre uma camada do problema; juntos formam o sistema")
    boxes = [
        ("PARA", "Tiago Forte", "Organizacao por nivel de acionabilidade",
         ["Projects: tem prazo", "Areas: continuas", "Resources: tematicos", "Archive: concluido"], TEAL),
        ("Zettelkasten", "Niklas Luhmann", "Notas atomicas linkadas",
         ["Uma ideia por nota", "Voz do aluno", "Links bidirecionais", "Crescimento organico"], NAVY),
        ("CODE", "Tiago Forte", "Workflow de processamento",
         ["Capture (inbox)", "Organize (PARA)", "Distill (Zettel)", "Express (drafts)"], GOLD),
    ]
    for i, (name, author, desc, bullets, color) in enumerate(boxes):
        x = Inches(0.5 + i*4.2)
        add_rect(s, x, Inches(1.3), Inches(4), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(4), Inches(0.9), color)
        add_text(s, x, Inches(1.4), Inches(4), Inches(0.5),
                 name, size=22, bold=True, color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(1.85), Inches(4), Inches(0.3),
                 author, size=11, color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.7 + i*4.2), Inches(2.4), Inches(3.6), Inches(0.5),
                 desc, size=13, bold=True, color=DARK)
        add_bullets(s, Inches(0.7 + i*4.2), Inches(3.0), Inches(3.6), Inches(3.5),
                    bullets, size=13, line_spacing=1.5)
    footer_bar(s, 3, TOTAL)


# ----------------------- 4. Plataforma -----------------------
def slide_plataforma():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Plataforma escolhida: Google Drive",
               "Trade-off analisado contra OneDrive")
    headers = ["Criterio", "Google Drive", "OneDrive", "Vencedor"]
    rows = [
        ("Espaco gratuito", "15 GB", "5 GB", "GDrive"),
        ("OCR em PDFs", "Nativo", "Limitado", "GDrive"),
        ("Integracao Scholar", "Library nativa", "Nao tem", "GDrive"),
        ("Google Colab (analise)", "Nativo", "Nao", "GDrive"),
        ("Word/Excel desktop", "Compativel", "Nativo", "OneDrive"),
        ("Versao gratis universidade", "Workspace EDU", "M365 EDU 1TB", "Depende"),
        ("Versionamento", "100 versoes", "25 versoes", "GDrive"),
    ]
    col_w = [Inches(3.5), Inches(3.5), Inches(3.5), Inches(2)]
    col_x = [Inches(0.5), Inches(4), Inches(7.5), Inches(11)]
    y = Inches(1.3)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.5), NAVY)
    for i, h in enumerate(headers):
        add_text(s, col_x[i], Inches(1.35), col_w[i], Inches(0.4),
                 h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        ry = Inches(1.85 + r*0.5)
        if r % 2 == 0:
            add_rect(s, Inches(0.5), ry, Inches(12.3), Inches(0.5), LIGHT)
        for i, val in enumerate(row):
            color = DARK
            bold = False
            if i == 3:
                color = GREEN if val == "GDrive" else (TEAL if val == "OneDrive" else GRAY)
                bold = True
            add_text(s, col_x[i], Inches(1.93 + r*0.5), col_w[i], Inches(0.4),
                     val, size=12, color=color, bold=bold, align=PP_ALIGN.CENTER)
    add_rect(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.4), TEAL)
    add_text(s, Inches(0.7), Inches(5.7), Inches(12), Inches(0.4),
             "Decisao: Google Drive como default",
             size=15, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(6.1), Inches(12), Inches(0.9),
             "Migre para OneDrive APENAS se sua universidade oferecer M365 EDU com 1TB e voce trabalhar intensivo no Word desktop + EndNote/Mendeley.",
             size=12, color=WHITE)
    footer_bar(s, 4, TOTAL)


# ----------------------- 5. Estrutura de pastas -----------------------
def slide_estrutura():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Estrutura de pastas no Google Drive",
               "Numeracao 00-99 garante ordem visual estavel")
    tree = [
        ("Mestrado-{Tema}/", 0, NAVY, True),
        ("00_Inbox/", 1, RED, False),
        ("01_Projects/  (P do PARA: tem prazo)", 1, TEAL, True),
        ("2026-Dissertacao/", 2, GRAY, False),
        ("2026-Q3-Artigo-Periodico/", 2, GRAY, False),
        ("2026-Q4-Qualificacao/", 2, GRAY, False),
        ("02_Areas/  (A do PARA: continuas)", 1, TEAL, True),
        ("Linha-Pesquisa/", 2, GRAY, False),
        ("Orientacao/  | Disciplinas/  | Carreira-Academica/", 2, GRAY, False),
        ("03_Resources/  (R do PARA: tematicos)", 1, TEAL, True),
        ("Bibliografia/ (PDFs/, Fichamentos/, biblio.bib)", 2, GRAY, False),
        ("Notas-Atomicas/  <-- nucleo Zettelkasten", 2, GOLD, True),
        ("Templates/  | Metodos/  | Glossario/", 2, GRAY, False),
        ("04_Archive/  (concluido ou inativo)", 1, TEAL, False),
        ("99_Sandbox/  (rascunhos descartaveis)", 1, GRAY, False),
    ]
    y = 1.3
    for label, indent, color, bold in tree:
        x = Inches(0.7 + indent * 0.4)
        add_text(s, x, Inches(y), Inches(11), Inches(0.32),
                 ("|-- " if indent > 0 else "") + label,
                 size=14, bold=bold, color=color, font="Consolas")
        y += 0.32
    add_rect(s, Inches(8.5), Inches(5.5), Inches(4.3), Inches(1.5), LIGHT)
    add_text(s, Inches(8.7), Inches(5.55), Inches(4), Inches(0.4),
             "Por que esta numeracao?",
             size=12, bold=True, color=NAVY)
    add_text(s, Inches(8.7), Inches(5.85), Inches(4), Inches(1.1),
             "Drive ordena lexicograficamente. 00 fica no topo (lembra de processar). 99 vai para o fundo (isola experimentos).",
             size=10, color=DARK)
    footer_bar(s, 5, TOTAL)


# ----------------------- 6. Convencao de nomenclatura -----------------------
def slide_nomenclatura():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Convencao de nomenclatura",
               "YYYYMMDD-tipo-slug-em-kebab-case.ext")
    add_rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0), NAVY)
    add_text(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.7),
             "20260507-lit-foucault-vigiar-punir.md",
             size=24, bold=True, color=GOLD, font="Consolas",
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(1.95), Inches(12), Inches(0.4),
             "data ISO ordenavel    tipo    slug em kebab-case    extensao",
             size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    types = [
        ("lit", "Fichamento de literatura"),
        ("nota", "Nota atomica permanente (Zettel)"),
        ("draft", "Rascunho de texto autoral"),
        ("dado", "Dataset bruto/tratado"),
        ("fig", "Figura, diagrama"),
        ("tab", "Tabela"),
        ("slides", "Apresentacao"),
        ("audio", "Audio (entrevista, voz do aluno)"),
        ("transcript", "Transcricao"),
        ("analise", "Output de analise"),
        ("peer", "Peer review"),
        ("meeting", "Ata de orientacao"),
    ]
    cols = 3
    for i, (prefix, desc) in enumerate(types):
        col = i % cols
        row = i // cols
        x = Inches(0.5 + col * 4.2)
        y = Inches(2.7 + row * 1.0)
        add_rect(s, x, y, Inches(4), Inches(0.85), LIGHT)
        add_rect(s, x, y, Inches(1.0), Inches(0.85), TEAL)
        add_text(s, x, y + Inches(0.15), Inches(1.0), Inches(0.55),
                 prefix, size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, x + Inches(1.1), y + Inches(0.2), Inches(2.8), Inches(0.5),
                 desc, size=11, color=DARK)
    footer_bar(s, 6, TOTAL)


# ----------------------- 7. CODE workflow -----------------------
def slide_code():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Workflow CODE adaptado",
               "Como o conteudo flui pelas pastas e pelos agentes")
    stages = [
        ("CAPTURE", "00_Inbox/", "Tudo cai aqui sem julgamento.\nPDF, foto, audio, snippet.", RED, "academic-researcher"),
        ("ORGANIZE", "03_Resources/", "Move para PARA.\nPDFs, biblio.bib, etc.", TEAL, "citation-manager"),
        ("DISTILL", "Notas-Atomicas/permanent/", "Uma ideia por nota.\nLinks bidirecionais.", GOLD, "knowledge-architect"),
        ("EXPRESS", "01_Projects/{atual}/drafts/", "Texto autoral.\nValidate > Humanize > Peer.", NAVY, "academic-writer + voice-humanizer + peer-reviewer"),
    ]
    for i, (name, folder, desc, color, agent) in enumerate(stages):
        x = Inches(0.4 + i * 3.25)
        add_rect(s, x, Inches(1.3), Inches(3), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(3), Inches(0.7), color)
        add_text(s, x, Inches(1.4), Inches(3), Inches(0.5),
                 name, size=18, bold=True,
                 color=WHITE if color != GOLD else DARK, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.15), Inches(2.2), Inches(2.7), Inches(0.4),
                 folder, size=11, bold=True, color=NAVY, font="Consolas")
        add_text(s, x + Inches(0.15), Inches(2.7), Inches(2.7), Inches(2),
                 desc, size=12, color=DARK)
        add_rect(s, x + Inches(0.15), Inches(5.0), Inches(2.7), Inches(0.05), color)
        add_text(s, x + Inches(0.15), Inches(5.1), Inches(2.7), Inches(0.4),
                 "AGENTE PRINCIPAL", size=9, bold=True, color=GRAY)
        add_text(s, x + Inches(0.15), Inches(5.4), Inches(2.7), Inches(1.3),
                 agent, size=11, color=DARK, bold=True)
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + Inches(3.0), Inches(3.7),
                                        Inches(0.3), Inches(0.4))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GOLD
            arrow.line.fill.background()
    footer_bar(s, 7, TOTAL)


# ----------------------- 8. Squad - mapa -----------------------
def slide_squad_mapa():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Squad de 8 agentes - mapa de fluxo",
               "Cada agente atua sobre pastas especificas do Google Drive")
    add_rect(s, Inches(5.5), Inches(1.3), Inches(2.3), Inches(0.7), NAVY)
    add_text(s, Inches(5.5), Inches(1.45), Inches(2.3), Inches(0.4),
             "VOCE", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(5.5), Inches(1.7), Inches(2.3), Inches(0.3),
             "(mestrando)", size=10, color=LIGHT, align=PP_ALIGN.CENTER)
    agents = [
        ("methodology-advisor", 0.8, 2.6, TEAL, "Desenho metodologico"),
        ("academic-researcher", 3.5, 2.6, TEAL, "Curadoria de fontes"),
        ("source-validator", 6.2, 2.6, RED, "Anti-alucinacao"),
        ("knowledge-architect", 8.9, 2.6, GOLD, "Zettelkasten"),
        ("academic-writer", 0.8, 5.0, NAVY, "Redacao ABNT/APA"),
        ("citation-manager", 3.5, 5.0, GRAY, "Bibliografia .bib"),
        ("voice-humanizer", 6.2, 5.0, GOLD, "Voz autoral"),
        ("peer-reviewer", 8.9, 5.0, RED, "Banca simulada"),
    ]
    for name, x, y, color, role in agents:
        add_rect(s, Inches(x), Inches(y), Inches(2.5), Inches(0.45), color)
        add_text(s, Inches(x), Inches(y + 0.05), Inches(2.5), Inches(0.4),
                 name, size=11, bold=True,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, Inches(x), Inches(y + 0.5), Inches(2.5), Inches(0.4),
                 role, size=10, color=DARK, align=PP_ALIGN.CENTER)
    footer_bar(s, 8, TOTAL)


# ----------------------- 9-12. Cada agente x ambiente -----------------------
def slide_agente_grupo1():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Como cada agente atua no Drive (1/4)",
               "Pesquisa: methodology-advisor + academic-researcher")
    agents_data = [
        {
            "name": "methodology-advisor",
            "color": TEAL,
            "le": ["02_Areas/Orientacao/", "01_Projects/{atual}/README.md", "memory/project_thesis.md"],
            "escreve": ["01_Projects/{atual}/metodologia.md", "memory/project_methodology.md"],
            "drive": "Usa Google Docs para metodologia se quiser comentarios do orientador. Mantem `memory/` versionado.",
            "trigger": "Inicio de cada novo projeto (artigo, capitulo, dissertacao)",
        },
        {
            "name": "academic-researcher",
            "color": TEAL,
            "le": ["memory/project_thesis.md", "memory/project_keywords.md"],
            "escreve": ["03_Resources/Bibliografia/PDFs/{citekey}.pdf", "01_Projects/{atual}/research/{data}-{slug}/sources.md"],
            "drive": "Aproveita botao Scholar > Save to Drive. Ao encontrar DOI, aciona citation-manager.",
            "trigger": "Revisao de literatura, busca dirigida, atualizacao do estado da arte",
        },
    ]
    for i, a in enumerate(agents_data):
        x = Inches(0.4 + i * 6.4)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(0.5), a["color"])
        add_text(s, x, Inches(1.38), Inches(6.3), Inches(0.4),
                 a["name"], size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(1.95), Inches(6), Inches(0.3),
                 "LE de:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["le"]):
            add_text(s, x + Inches(0.3), Inches(2.25 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(3.3), Inches(6), Inches(0.3),
                 "ESCREVE em:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["escreve"]):
            add_text(s, x + Inches(0.3), Inches(3.6 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(4.6), Inches(6), Inches(0.3),
                 "INTEGRACAO Drive:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(4.9), Inches(6), Inches(1.0),
                 a["drive"], size=10, color=DARK)
        add_text(s, x + Inches(0.2), Inches(5.85), Inches(6), Inches(0.3),
                 "TRIGGER:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(6.15), Inches(6), Inches(0.6),
                 a["trigger"], size=10, color=DARK)
    footer_bar(s, 9, TOTAL)


def slide_agente_grupo2():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Como cada agente atua no Drive (2/4)",
               "Validacao e organizacao: source-validator + knowledge-architect")
    agents_data = [
        {
            "name": "source-validator",
            "color": RED,
            "le": ["01_Projects/{atual}/drafts/*.md", "03_Resources/Bibliografia/biblio.bib"],
            "escreve": ["01_Projects/{atual}/peer-review/{data}-validacao-{slug}.md"],
            "drive": "Outputs ficam dentro do projeto auditado. Nunca em recursos compartilhados.",
            "trigger": "APOS qualquer output de IA. ANTES de qualquer submissao.",
        },
        {
            "name": "knowledge-architect",
            "color": GOLD,
            "le": ["PDFs em Bibliografia/PDFs/", "Fichamentos brutos", "Audios dictados"],
            "escreve": ["03_Resources/Bibliografia/Fichamentos/{citekey}.md",
                        "03_Resources/Notas-Atomicas/permanent/{ts}-{slug}.md",
                        "03_Resources/Notas-Atomicas/index/moc-{tema}.md"],
            "drive": "Notas-Atomicas/ e o nucleo do segundo cerebro. Backup adicional em git privado recomendado.",
            "trigger": "A cada leitura concluida, insight, ou no rito quinzenal de promocao raw>refined>connected.",
        },
    ]
    for i, a in enumerate(agents_data):
        x = Inches(0.4 + i * 6.4)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(0.5), a["color"])
        add_text(s, x, Inches(1.38), Inches(6.3), Inches(0.4),
                 a["name"], size=15, bold=True,
                 color=WHITE if a["color"] != GOLD else DARK,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(1.95), Inches(6), Inches(0.3),
                 "LE de:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["le"]):
            add_text(s, x + Inches(0.3), Inches(2.25 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(3.3), Inches(6), Inches(0.3),
                 "ESCREVE em:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["escreve"]):
            add_text(s, x + Inches(0.3), Inches(3.6 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(4.85), Inches(6), Inches(0.3),
                 "INTEGRACAO Drive:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(5.15), Inches(6), Inches(1.0),
                 a["drive"], size=10, color=DARK)
        add_text(s, x + Inches(0.2), Inches(6.0), Inches(6), Inches(0.3),
                 "TRIGGER:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(6.3), Inches(6), Inches(0.6),
                 a["trigger"], size=10, color=DARK)
    footer_bar(s, 10, TOTAL)


def slide_agente_grupo3():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Como cada agente atua no Drive (3/4)",
               "Producao: academic-writer + citation-manager")
    agents_data = [
        {
            "name": "academic-writer",
            "color": NAVY,
            "le": ["Notas-Atomicas/permanent/", "biblio.bib", "memory/user_writing_voice.md"],
            "escreve": ["01_Projects/{atual}/drafts/{capitulo}.md"],
            "drive": "Ao terminar capitulo, exporta .md > .docx (Pandoc) para enviar ao orientador via Google Docs.",
            "trigger": "Sprint de redacao. Sempre apos o aluno selecionar 8-12 notas conectadas.",
        },
        {
            "name": "citation-manager",
            "color": GRAY,
            "le": ["Listas de referencias em qualquer formato", "DOIs"],
            "escreve": ["03_Resources/Bibliografia/biblio.bib (UNICO no Drive)",
                        "01_Projects/{atual}/referencias-{norma}.md"],
            "drive": "biblio.bib e arquivo unico compartilhado entre todos os capitulos. Versionado pelo Drive.",
            "trigger": "Toda nova fonte. Antes de cada submissao. Conversao entre normas.",
        },
    ]
    for i, a in enumerate(agents_data):
        x = Inches(0.4 + i * 6.4)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(0.5), a["color"])
        add_text(s, x, Inches(1.38), Inches(6.3), Inches(0.4),
                 a["name"], size=15, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(1.95), Inches(6), Inches(0.3),
                 "LE de:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["le"]):
            add_text(s, x + Inches(0.3), Inches(2.25 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(3.3), Inches(6), Inches(0.3),
                 "ESCREVE em:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["escreve"]):
            add_text(s, x + Inches(0.3), Inches(3.6 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(4.6), Inches(6), Inches(0.3),
                 "INTEGRACAO Drive:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(4.9), Inches(6), Inches(1.0),
                 a["drive"], size=10, color=DARK)
        add_text(s, x + Inches(0.2), Inches(5.85), Inches(6), Inches(0.3),
                 "TRIGGER:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(6.15), Inches(6), Inches(0.6),
                 a["trigger"], size=10, color=DARK)
    footer_bar(s, 11, TOTAL)


def slide_agente_grupo4():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Como cada agente atua no Drive (4/4)",
               "Polimento: voice-humanizer + peer-reviewer")
    agents_data = [
        {
            "name": "voice-humanizer",
            "color": GOLD,
            "le": ["memory/voice_questionnaire.md", "memory/user_writing_voice.md", "drafts a humanizar"],
            "escreve": ["memory/user_writing_voice.md (ATUALIZA)",
                        "drafts reescritos em colaboracao",
                        "memory/voice_calibration_log.md"],
            "drive": "Aplica QUESTIONARIO de calibracao na primeira sessao. Aprende voz do aluno e atualiza perfil ao longo do tempo.",
            "trigger": "Antes de qualquer submissao. Re-calibracao trimestral.",
        },
        {
            "name": "peer-reviewer",
            "color": RED,
            "le": ["drafts completos", "metodologia.md", "biblio.bib", "MOCs"],
            "escreve": ["01_Projects/{atual}/peer-review/{data}-banca-{slug}.md"],
            "drive": "Gera versao .docx para imprimir e levar a reuniao com orientador. Mantem banco de perguntas em memory/banca_questions.md.",
            "trigger": "ANTES de qualificacao, defesa ou submissao a revista.",
        },
    ]
    for i, a in enumerate(agents_data):
        x = Inches(0.4 + i * 6.4)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(6.3), Inches(0.5), a["color"])
        add_text(s, x, Inches(1.38), Inches(6.3), Inches(0.4),
                 a["name"], size=15, bold=True,
                 color=WHITE if a["color"] != GOLD else DARK,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(1.95), Inches(6), Inches(0.3),
                 "LE de:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["le"]):
            add_text(s, x + Inches(0.3), Inches(2.25 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(3.3), Inches(6), Inches(0.3),
                 "ESCREVE em:", size=11, bold=True, color=NAVY)
        for j, item in enumerate(a["escreve"]):
            add_text(s, x + Inches(0.3), Inches(3.6 + j*0.3), Inches(6), Inches(0.3),
                     "- " + item, size=10, color=DARK, font="Consolas")
        add_text(s, x + Inches(0.2), Inches(4.85), Inches(6), Inches(0.3),
                 "INTEGRACAO Drive:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(5.15), Inches(6), Inches(1.0),
                 a["drive"], size=10, color=DARK)
        add_text(s, x + Inches(0.2), Inches(6.05), Inches(6), Inches(0.3),
                 "TRIGGER:", size=11, bold=True, color=NAVY)
        add_text(s, x + Inches(0.3), Inches(6.35), Inches(6), Inches(0.6),
                 a["trigger"], size=10, color=DARK)
    footer_bar(s, 12, TOTAL)


# ----------------------- 13. Voice-humanizer questionario -----------------------
def slide_voice_questionario():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "voice-humanizer: questionario de calibracao",
               "21 perguntas em 5 blocos. Aplicado na primeira sessao.")
    blocks = [
        ("Bloco 1", "Identidade academica", "5 min",
         "Area, nivel de formalidade, eu/nos/impessoal, autores admirados/rejeitados",
         TEAL),
        ("Bloco 2", "Repertorio linguistico", "10 min",
         "5 palavras que voce usa, 5 que rejeita, frase curta/longa, verbo padrao de posicionamento",
         NAVY),
        ("Bloco 3", "Habitos de escrita", "10 min",
         "Direto/oral/manual, como reage ao bloqueio, tempo em revisao, sensacao ao reler",
         GOLD),
        ("Bloco 4", "Posicionamento e ethos", "5 min",
         "Conciliador/provocador/didatico/sintetico, como aborda autor citado, tema com urgencia",
         RED),
        ("Bloco 5", "Materiais para analise", "5 min",
         "3 textos seus bons + 1 odiado + 1 informal (e-mail) para extracao de padroes",
         GRAY),
    ]
    for i, (b, t, time, desc, color) in enumerate(blocks):
        y = Inches(1.3 + i * 1.05)
        add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.95), LIGHT)
        add_rect(s, Inches(0.5), y, Inches(1.4), Inches(0.95), color)
        add_text(s, Inches(0.5), y + Inches(0.2), Inches(1.4), Inches(0.4),
                 b, size=14, bold=True,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(0.5), y + Inches(0.5), Inches(1.4), Inches(0.3),
                 time, size=10,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(2.0), y + Inches(0.1), Inches(10.7), Inches(0.4),
                 t, size=14, bold=True, color=NAVY)
        add_text(s, Inches(2.0), y + Inches(0.45), Inches(10.7), Inches(0.5),
                 desc, size=11, color=DARK)
    add_rect(s, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5), GOLD)
    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.3),
             "Saida: memory/voice_questionnaire.md (bruto) + memory/user_writing_voice.md (perfil destilado)",
             size=12, bold=True, color=DARK)
    footer_bar(s, 13, TOTAL)


# ----------------------- 14. Pesquisa rapida -----------------------
def slide_pesquisa_rapida():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Como achar qualquer coisa em 5 segundos",
               "5 estrategias de busca cobrindo todos os cenarios")
    strategies = [
        ("Drive Web", "Interface oficial",
         ["type:pdf foucault", "before:2026-01-01 type:document", "OCR busca dentro do conteudo do PDF"], TEAL),
        ("Windows Search", "Drive for Desktop",
         ["Ctrl+F no Explorer", "Filtro por nome ISO ordenavel", "Pre-visualizacao instantanea"], NAVY),
        ("Busca semantica", "via knowledge-architect",
         ["'Liste notas que conectam Foucault e Deleuze'", "Agente busca por relacao semantica", "Cobre o que tag nao captura"], GOLD),
        ("Busca por tag (grep)", "Frontmatter YAML",
         ["grep -lr 'tags:.*foucault' Notas-Atomicas/", "Filtro por status: connected | refined", "Listagem de notas por projeto"], RED),
        ("MOCs (Map of Content)", "Curado manualmente",
         ["Notas-Atomicas/index/moc-{tema}.md", "Portao de entrada por tema", "Comeca aqui antes de buscar"], GRAY),
    ]
    for i, (name, sub, items, color) in enumerate(strategies):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.2)
        y = Inches(1.3 + row * 2.8)
        add_rect(s, x, y, Inches(4), Inches(2.6), LIGHT)
        add_rect(s, x, y, Inches(4), Inches(0.55), color)
        add_text(s, x, y + Inches(0.05), Inches(4), Inches(0.3),
                 name, size=14, bold=True,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.32), Inches(4), Inches(0.25),
                 sub, size=10,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_bullets(s, x + Inches(0.15), y + Inches(0.7), Inches(3.7), Inches(1.9),
                    items, size=10, line_spacing=1.3)
    footer_bar(s, 14, TOTAL)


# ----------------------- 15. Cronograma -----------------------
def slide_cronograma():
    s = prs.slides.add_slide(BLANK)
    header_bar(s, "Cronograma de adocao - 4 semanas",
               "De zero ao primeiro fluxo CODE completo")
    weeks = [
        ("Semana 1", "Estrutura", TEAL,
         ["Criar pasta Mestrado-{Tema} no Drive",
          "Replicar arvore PARA + Zettelkasten",
          "Instalar Drive for Desktop (sync local)",
          "Mover material disperso para 00_Inbox/"]),
        ("Semana 2", "Calibracao", NAVY,
         ["Preencher memory/project_thesis.md",
          "Preencher memory/citation_style.md",
          "APLICAR questionario do voice-humanizer",
          "Submeter 3 textos antigos para extracao de voz"]),
        ("Semana 3", "Primeiro CODE", GOLD,
         ["Capturar 5 fontes via academic-researcher",
          "Validar com source-validator",
          "Fichar com knowledge-architect",
          "Linkar 5 notas atomicas bidirecionalmente"]),
        ("Semana 4", "Primeiro Express", RED,
         ["Selecionar 8-12 notas conectadas",
          "academic-writer produz mini-secao 1500 palavras",
          "Ciclo: validate > humanize > peer-review",
          "Avaliar e ajustar o sistema"]),
    ]
    for i, (week, title, color, items) in enumerate(weeks):
        x = Inches(0.4 + i * 3.25)
        add_rect(s, x, Inches(1.3), Inches(3), Inches(5.5), LIGHT)
        add_rect(s, x, Inches(1.3), Inches(3), Inches(0.85), color)
        add_text(s, x, Inches(1.4), Inches(3), Inches(0.4),
                 week, size=15, bold=True,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_text(s, x, Inches(1.75), Inches(3), Inches(0.4),
                 title, size=18, bold=True,
                 color=WHITE if color != GOLD else DARK,
                 align=PP_ALIGN.CENTER)
        add_bullets(s, x + Inches(0.15), Inches(2.4), Inches(2.7), Inches(4.2),
                    items, size=11, line_spacing=1.4)
    footer_bar(s, 15, TOTAL)


# ----------------------- 16. Proximos passos -----------------------
def slide_proximos_passos():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, NAVY)
    add_rect(s, Inches(0), Inches(2.8), prs.slide_width, Inches(0.08), GOLD)
    add_text(s, Inches(0.7), Inches(1.0), Inches(12), Inches(0.8),
             "Proximos passos",
             size=40, bold=True, color=WHITE)
    add_text(s, Inches(0.7), Inches(1.9), Inches(12), Inches(0.5),
             "Para sair daqui com o sistema operando",
             size=18, color=GOLD)
    steps = [
        ("1", "Aprovar este plano", "Revisar pastas, nomenclatura, agentes - ajustar ao seu tema antes de implementar"),
        ("2", "Provisionar o Drive", "Criar Mestrado-{Tema}/ com toda a arvore PARA. Instalar Drive for Desktop."),
        ("3", "Aplicar questionario", "voice-humanizer aplica os 21 itens. Voce envia 3 textos antigos para calibracao."),
        ("4", "Rodar Sprint 1", "Capturar 5 fontes > validar > fichar > 5 notas atomicas. Validar fluxo end-to-end."),
        ("5", "Avaliar e iterar", "Apos primeira mini-secao escrita, ajustar prompts dos agentes que falharem."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(3.2 + i * 0.75)
        add_rect(s, Inches(0.7), y, Inches(0.6), Inches(0.6), GOLD)
        add_text(s, Inches(0.7), y + Inches(0.1), Inches(0.6), Inches(0.5),
                 num, size=22, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(s, Inches(1.5), y + Inches(0.05), Inches(11), Inches(0.35),
                 title, size=16, bold=True, color=WHITE)
        add_text(s, Inches(1.5), y + Inches(0.35), Inches(11), Inches(0.3),
                 desc, size=11, color=LIGHT)
    add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
             "Documentos: docs/curadoria-conteudo-plano.md  |  agents/  |  .claude/plan/agentes-mestrado-plano.md",
             size=10, color=GRAY, align=PP_ALIGN.CENTER, font="Consolas")


# ----------------------- build -----------------------
slide_capa()
slide_problema()
slide_visao()
slide_plataforma()
slide_estrutura()
slide_nomenclatura()
slide_code()
slide_squad_mapa()
slide_agente_grupo1()
slide_agente_grupo2()
slide_agente_grupo3()
slide_agente_grupo4()
slide_voice_questionario()
slide_pesquisa_rapida()
slide_cronograma()
slide_proximos_passos()

out_dir = Path(__file__).parent.parent / "docs"
out_dir.mkdir(exist_ok=True)
out_path = out_dir / "curadoria-conteudo-framework.pptx"
prs.save(str(out_path))
print(f"PPTX gerado: {out_path}")
print(f"Total de slides: {len(prs.slides)}")
